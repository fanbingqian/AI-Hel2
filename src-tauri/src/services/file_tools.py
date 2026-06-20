"""
File tools: extract text from documents and images for LLM summarization.
Usage:
  python file_tools.py extract <file_path> [--type pdf|docx|pptx|image]
Output: JSON {"text": "...", "pages": N, "error": null}
"""

import sys
import os
import json
import argparse
import base64
import io

sys.stdout.reconfigure(encoding="utf-8")


def extract_markitdown(path):
    """Try Microsoft MarkItDown for unified DOCX/PPTX/PDF/XLSX → Markdown.
    Returns Markdown text on success, None if not installed or conversion fails."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(path)
        if result and result.text_content and result.text_content.strip():
            return result.text_content.strip()
    except ImportError:
        pass
    except Exception:
        pass
    return None


def extract_pdf(path):
    """Extract text from PDF using MarkItDown (preferred), pdfplumber, or PyPDF2."""
    # Try MarkItDown first (best quality: Markdown output)
    md_text = extract_markitdown(path)
    if md_text:
        pages = md_text.count("\n\n") + 1
        return md_text, max(pages, 1), None

    # Fallback to pdfplumber / PyPDF2
    text_parts = []
    pages = 0
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
                pages += 1
        return "\n\n".join(text_parts), pages, None
    except ImportError:
        pass
    except Exception as e:
        pass  # fall through to PyPDF2

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
            pages += 1
        return "\n\n".join(text_parts), pages, None
    except ImportError:
        return "", 0, "No PDF library available. Install pdfplumber or PyPDF2."
    except Exception as e:
        return "", 0, f"PDF extract error: {e}"


def extract_docx(path):
    """Extract text from Word .docx using MarkItDown (preferred) or python-docx."""
    md_text = extract_markitdown(path)
    if md_text:
        lines = md_text.split("\n")
        return md_text, len([l for l in lines if l.strip()]), None
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs), len(paragraphs), None
    except ImportError:
        return "", 0, "No Word library available. Install python-docx."
    except Exception as e:
        return "", 0, f"DOCX extract error: {e}"


def extract_pptx(path):
    """Extract text from PowerPoint .pptx using MarkItDown (preferred) or python-pptx."""
    md_text = extract_markitdown(path)
    if md_text:
        lines = md_text.split("\n")
        return md_text, len([l for l in lines if l.strip()]), None
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if slide_parts:
                slides_text.append(f"## 幻灯片 {i + 1}\n" + "\n".join(slide_parts))
        return "\n\n".join(slides_text), len(slides_text), None
    except ImportError:
        return "", 0, "No PPTX library available. Install python-pptx."
    except Exception as e:
        return "", 0, f"PPTX extract error: {e}"


def extract_xlsx(path):
    """Extract text from Excel .xlsx using MarkItDown (preferred) or openpyxl."""
    md_text = extract_markitdown(path)
    if md_text:
        lines = md_text.split("\n")
        return md_text, len([l for l in lines if l.strip()]), None
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    rows_text.append(" | ".join(cells))
            if rows_text:
                parts.append(f"## 工作表: {sheet_name}\n" + "\n".join(rows_text))
        wb.close()
        return "\n\n".join(parts), len(parts), None
    except ImportError:
        return "", 0, "No Excel library available. Install openpyxl."
    except Exception as e:
        return "", 0, f"XLSX extract error: {e}"


def extract_image(path):
    """Extract text from image: OCR first, then base64 fallback for LLM."""
    # Try OCR first
    try:
        from PIL import Image
        img = Image.open(path)
        try:
            import pytesseract
            text = pytesseract.image_to_string(img, lang="chi_sim+eng")
            if text and text.strip():
                return text.strip(), 1, None
        except ImportError:
            pass
        except Exception:
            pass
    except Exception:
        pass

    # Fallback: base64 for LLM multimodal description
    b64, err = describe_image_base64(path)
    if err:
        return "", 0, err
    # Return base64 as "text" — caller's LLM handles multimodal
    return b64, 1, None


def describe_image_base64(path):
    """Read image file and return base64 data URL for multimodal LLM."""
    try:
        ext = os.path.splitext(path)[1].lower()
        mime_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
            ".svg": "image/svg+xml",
        }
        mime = mime_map.get(ext, "image/png")
        with open(path, "rb") as f:
            data = base64.b64encode(f.read()).decode("ascii")
        return f"data:{mime};base64,{data}", None
    except Exception as e:
        return "", str(e)


def main():
    parser = argparse.ArgumentParser(description="File tools for document text extraction")
    parser.add_argument("action", choices=["extract", "image_b64"])
    parser.add_argument("path", help="File path")
    parser.add_argument("--type", default=None, help="File type override")
    args = parser.parse_args()

    path = args.path
    if not os.path.exists(path):
        print(json.dumps({"text": "", "pages": 0, "error": f"File not found: {path}"}))
        sys.exit(1)

    if args.action == "image_b64":
        b64, error = describe_image_base64(path)
        if error:
            print(json.dumps({"text": "", "pages": 0, "error": error}))
            sys.exit(1)
        else:
            print(json.dumps({"text": b64, "pages": 1, "error": None}))
        return

    # Determine file type
    ext = os.path.splitext(path)[1].lower()
    file_type = args.type or ext.lstrip(".")

    extractors = {
        "pdf": extract_pdf,
        "docx": extract_docx,
        "pptx": extract_pptx,
        "xlsx": extract_xlsx,
        "png": extract_image, "jpg": extract_image, "jpeg": extract_image,
        "gif": extract_image, "webp": extract_image, "bmp": extract_image,
        "svg": extract_image, "ico": extract_image, "tiff": extract_image,
    }

    if file_type in extractors:
        text, pages, error = extractors[file_type](path)
        result = {"text": text, "pages": pages, "error": error}
    else:
        # Only attempt text read for known text-based extensions
        text_exts = {"txt", "csv", "md", "json", "xml", "html", "css", "js", "py", "rs", "ts", "yaml", "yml", "log", "cfg", "ini", "toml"}
        if ext.lstrip(".") in text_exts:
            try:
                with open(path, "r", encoding="utf-8") as f:
                    text = f.read()
                result = {"text": text, "pages": 1, "error": None}
            except Exception:
                result = {"text": "", "pages": 0, "error": f"Cannot read file: {file_type}"}
        else:
            result = {"text": "", "pages": 0, "error": f"Unsupported file type: {file_type}"}

    sys.stdout.write(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
