"""HEIMDALL Web Console API — REST + SSE routes.

FastAPI application with routes for all five V3.0 human views plus
SSE streaming chat backed by the existing AIAgent + HeimdallManager.
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
import re
import sys
import time
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, File, HTTPException, Query, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles

# Ensure the project root is on sys.path
_project_root = Path(__file__).resolve().parent.parent.parent
if str(_project_root) not in sys.path:
    sys.path.insert(0, str(_project_root))

# Load .env from ~/.hermes/.env (or HERMES_HOME/.env) so API keys are available
from dotenv import load_dotenv as _load_dotenv
from hermes_constants import get_hermes_home as _get_hermes_home
_env_path = _get_hermes_home() / ".env"
if _env_path.exists():
    _load_dotenv(_env_path)

logger = logging.getLogger(__name__)

app = FastAPI(title="HEIMDALL Console", version="0.2.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])


@app.on_event("startup")
async def _start_pending_retry():
    """Start background thread for retrying failed pending_sync entries."""
    global _pending_retry_running, _pending_retry_thread
    import threading
    if not _pending_retry_running:
        _pending_retry_running = True
        _pending_retry_thread = threading.Thread(target=_pending_retry_loop, daemon=True)
        _pending_retry_thread.start()
        logger.info("Pending sync retry loop started")


@app.on_event("shutdown")
async def _stop_pending_retry():
    """Stop the pending_sync retry thread."""
    global _pending_retry_running
    _pending_retry_running = False

# ---------------------------------------------------------------------------
# Global state (lazy init on first request)
# ---------------------------------------------------------------------------

_heimdall = None
_agent = None
_chat_history: list[dict] = []  # persisted across the server session
_pending_retry_thread = None
_pending_retry_running = False


def _pending_retry_loop():
    """Background thread: retry kr_pending_sync entries every 30s."""
    import time as _time
    from datetime import datetime as _dt, timedelta
    global _pending_retry_running
    logger.info("Pending sync retry thread started")
    while _pending_retry_running:
        try:
            hm = _heimdall
            if hm and hm.provider.store and hm.provider.store._conn:
                store = hm.provider.store
                conn = store._conn
                now = _dt.now().isoformat()
                rows = conn.execute(
                    "SELECT * FROM kr_pending_sync WHERE next_retry_at <= ? AND retry_count < 5",
                    (now,),
                ).fetchall()
                for row in rows:
                    try:
                        payload = json.loads(row["payload_json"])
                        namespace = row["namespace"] or "general"
                        if row["payload_type"] == "wiki":
                            result = hm.provider.extractor.extract_from_text(
                                payload.get("content", ""), session_id="pending-retry", namespace=namespace
                            )
                        elif row["payload_type"] == "canvas":
                            # Actually re-process canvas nodes/edges (was a no-op stub)
                            nodes = payload.get("data", {}).get("nodes", [])
                            edges_data = payload.get("data", {}).get("edges", [])
                            entity_count = 0
                            relation_count = 0
                            for node in nodes:
                                node_type = node.get("_hermes", {}).get("entityType", "concept")
                                store.upsert_entity_v2(
                                    name=node.get("label", node.get("id", "")),
                                    entity_type=node_type,
                                    properties={"canvas_id": node.get("id", ""),
                                                "description": node.get("description", "")},
                                    namespace=namespace,
                                )
                                entity_count += 1
                            for edge in edges_data:
                                src = edge.get("source", "")
                                tgt = edge.get("target", "")
                                if src and tgt:
                                    src_entity = store.get_entity_v2_by_name(src)
                                    tgt_entity = store.get_entity_v2_by_name(tgt)
                                    if not src_entity:
                                        eid = store.upsert_entity_v2(name=src, namespace=namespace)
                                        src_entity = store.get_entity_v2(eid)
                                    if not tgt_entity:
                                        eid = store.upsert_entity_v2(name=tgt, namespace=namespace)
                                        tgt_entity = store.get_entity_v2(eid)
                                    if src_entity and tgt_entity:
                                        store.add_relation(
                                            source_id=src_entity["entity_id"],
                                            target_id=tgt_entity["entity_id"],
                                            rel_type=edge.get("label", "relates_to"),
                                            namespace=namespace,
                                        )
                                        relation_count += 1
                            result = {"synced": True, "entity_count": entity_count,
                                      "relation_count": relation_count}
                        else:
                            conn.execute("DELETE FROM kr_pending_sync WHERE id = ?", (row["id"],))
                            continue

                        conn.execute("DELETE FROM kr_pending_sync WHERE id = ?", (row["id"],))
                        logger.info("Retried pending sync %s successfully", row["id"])
                    except Exception as e:
                        retry_count = (row["retry_count"] or 0) + 1
                        delay = min(2 ** retry_count * 60, 3600)
                        next_retry = (_dt.now() + timedelta(seconds=delay)).isoformat()
                        conn.execute(
                            "UPDATE kr_pending_sync SET retry_count = ?, last_error = ?, "
                            "next_retry_at = ? WHERE id = ?",
                            (retry_count, str(e)[:500], next_retry, row["id"]),
                        )
                        if retry_count >= 5:
                            logger.warning("Pending sync %s marked dead after %d retries",
                                          row["id"], retry_count)

                # Clean up dead rows older than 7 days (prevents unbounded table growth)
                try:
                    dead_cutoff = (_dt.now() - timedelta(days=7)).isoformat()
                    deleted = conn.execute(
                        "DELETE FROM kr_pending_sync WHERE retry_count >= 5 AND next_retry_at < ?",
                        (dead_cutoff,),
                    )
                    conn.commit()
                    if deleted.rowcount > 0:
                        logger.info("Cleaned %d dead pending sync rows", deleted.rowcount)
                except Exception:
                    pass
        except Exception:
            pass
        _time.sleep(30)


def _make_extract_fn():
    """Build an extract_fn callable for LLM-assisted entity extraction.

    Returns a callable (system_prompt, user_prompt) -> str that uses the
    same provider/model as the main agent, or None if LLM is unavailable.
    """
    try:
        provider, model = _get_model_config()
        from agent.auxiliary_client import call_llm as _call_llm

        def _extract(system_prompt: str, user_prompt: str) -> str:
            resp = _call_llm(
                provider=provider,
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=0.2,
                max_tokens=2000,
                timeout=25,
                extra_body={"reasoning_effort": "low"},
            )
            content = resp.choices[0].message.content
            # DeepSeek reasoning models may put output in reasoning_content when
            # content is empty (all tokens consumed by reasoning)
            if not content:
                rc = getattr(resp.choices[0].message, "reasoning_content", None)
                if rc:
                    # Extract JSON from reasoning content
                    import re as _re
                    m = _re.search(r'\{.*\}', rc, _re.DOTALL)
                    if m:
                        content = m.group(0)
            return content or ""

        # Smoke-test the extract function to catch misconfiguration early
        _extract("Say 'ok'", "ok test")
        return _extract
    except Exception as exc:
        logger.warning("LLM extraction unavailable, using rule-based fallback: %s", exc)
        return None


def get_heimdall():
    global _heimdall
    if _heimdall is None:
        from heimdall.manager import HeimdallManager
        _sid = f"web-console-{uuid.uuid4().hex[:8]}"
        _extract_fn = _make_extract_fn()
        _heimdall = HeimdallManager(session_id=_sid, extract_fn=_extract_fn)
        _heimdall.initialize(session_id=_sid)
        logger.info("HeimdallManager initialized, session=%s, llm_extraction=%s",
                     _sid, _extract_fn is not None)
        # Initialize Knowledge Ring DailyReportEngine
        from heimdall.views.daily import set_daily_engine
        set_daily_engine(_heimdall.provider.store)
        _maybe_auto_snapshot()
    return _heimdall


def _get_model_config():
    """Read model and provider from ~/.hermes/config.yaml with env overrides."""
    provider = os.getenv("HERMES_INFERENCE_PROVIDER", "")
    model = os.getenv("LLM_MODEL", "")
    try:
        from hermes_cli.config import load_config
        cfg = load_config()
        if not provider:
            provider = cfg.get("model", {}).get("provider", "deepseek")
        if not model:
            model = cfg.get("model", {}).get("default", "deepseek-chat")
    except Exception:
        if not provider:
            provider = "deepseek"
        if not model:
            model = "deepseek-chat"
    return provider, model


def get_agent():
    global _agent
    if _agent is None:
        from run_agent import AIAgent
        provider, model = _get_model_config()
        _agent = AIAgent(
            quiet_mode=True,
            skip_context_files=True,
            provider=provider,
            model=model,
        )
        _agent._heimdall = get_heimdall()

        # Pre-build system prompt so chat is ready
        _agent._cached_system_prompt = _agent._build_system_prompt()
        logger.info("AIAgent initialized, model=%s", _agent.model)
    return _agent


# ---------------------------------------------------------------------------
# Persona Snapshot — time-series for drift curve
# ---------------------------------------------------------------------------

class _PersonaSnapshot:
    """Manages persona snapshots at ~/.hermes/persona_snapshots.json."""

    def __init__(self):
        self._path = _get_hermes_home() / "persona_snapshots.json"

    def _read(self) -> dict:
        if not self._path.exists():
            return {"snapshots": []}
        try:
            return json.loads(self._path.read_text("utf-8"))
        except Exception:
            return {"snapshots": []}

    def _write(self, data: dict):
        self._path.parent.mkdir(parents=True, exist_ok=True)
        self._path.write_text(json.dumps(data, ensure_ascii=False, indent=2), "utf-8")

    def capture(self, persona_manager) -> dict:
        """Capture current persona state, append to snapshots file."""
        data = self._read()
        now = time.time()
        today = datetime.fromtimestamp(now).strftime("%Y-%m-%d")

        # Keep only one snapshot per day
        data["snapshots"] = [
            s for s in data["snapshots"]
            if datetime.fromtimestamp(s["timestamp"]).strftime("%Y-%m-%d") != today
        ]

        drift = persona_manager.check_drift()
        snapshot = {
            "id": uuid.uuid4().hex[:12],
            "timestamp": now,
            "date": today,
            "core_self": persona_manager.get_core_self(),
            "external_persona": persona_manager.get_external_persona(),
            "user_profile": persona_manager.get_user_profile(),
            "social_anchors": persona_manager.get_social_anchors(),
            "drift_check": drift,
        }
        data["snapshots"].append(snapshot)
        self._write(data)
        logger.info("Persona snapshot captured, id=%s", snapshot["id"])
        return snapshot

    def get_drift_data(self, persona_manager) -> dict:
        """Return drift time-series + current drift check."""
        data = self._read()
        current_drift = persona_manager.check_drift()
        series = []
        for s in data["snapshots"]:
            d = s.get("drift_check", {})
            series.append({
                "date": s.get("date", ""),
                "timestamp": s.get("timestamp", 0),
                "core_value_count": len(s.get("core_self", [])),
                "drift_score": d.get("cosine_similarity", 1.0),
                "drifted": d.get("drifted", False),
                "needs_review": d.get("needs_review", False),
            })
        return {"snapshots": series, "current_drift": current_drift}

    def last_snapshot_time(self) -> float:
        data = self._read()
        if data["snapshots"]:
            return data["snapshots"][-1]["timestamp"]
        return 0


_persona_snapshot = _PersonaSnapshot()


def _maybe_auto_snapshot():
    """Capture a persona snapshot if none exists or > 1hr since last."""
    try:
        hm = get_heimdall()
        pm = hm.provider.persona
        if not pm:
            return
        last_ts = _persona_snapshot.last_snapshot_time()
        if last_ts == 0 or (time.time() - last_ts) > 3600:
            _persona_snapshot.capture(pm)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Static files
# ---------------------------------------------------------------------------

_static_dir = Path(__file__).resolve().parent / "static"


@app.get("/")
async def index():
    html = (_static_dir / "index.html").read_text(encoding="utf-8")
    return HTMLResponse(html)


@app.get("/mobile")
async def mobile_index():
    html = (_static_dir / "mobile" / "index.html").read_text(encoding="utf-8")
    return HTMLResponse(html)


# Mount static files for assets (CSS/JS/images) under /static/
if _static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(_static_dir)), name="static")


# ---------------------------------------------------------------------------
# Stats + Elevator
# ---------------------------------------------------------------------------

@app.get("/api/stats")
async def api_stats():
    hm = get_heimdall()
    stats = hm.get_stats()
    gs = hm.provider.store
    sg = hm.provider.social_graph

    return {
        "entities": stats.get("entities", 0),
        "is_first_run": stats.get("is_first_run", False),
        "entity_types": _get_entity_type_counts(gs),
        "social_graph": sg.get_stats() if sg else {},
        "elevator": _get_elevator_stats(hm),
    }


@app.get("/api/elevator/classify")
async def api_classify(q: str = ""):
    hm = get_heimdall()
    level = hm.provider.classify_query(q) if q else None
    level_name = level.name if level and hasattr(level, "name") else "L1"
    warning = hm.provider.get_elevator_warning(level, "") if level and level_name != "L1" else None
    return {"query": q, "level": level_name, "warning": warning}


# ---------------------------------------------------------------------------
# Persona (four-layer boundaries)
# ---------------------------------------------------------------------------

def _persona_content(val):
    """Normalize persona layer value to a string."""
    if isinstance(val, list):
        return "\n".join(val)
    return val or ""


@app.get("/api/persona")
async def api_persona():
    hm = get_heimdall()
    pm = hm.provider.persona
    if not pm:
        return {"layers": [], "empty": True}
    layers = [
        {"name": "core_self", "label": "核心自我", "icon": "🛡️", "content": _persona_content(pm.get_core_self()), "editable": False, "immune": True},
        {"name": "external_persona", "label": "对外人格", "icon": "🎭", "content": _persona_content(pm.get_external_persona()), "editable": False},
        {"name": "user_profile", "label": "用户画像", "icon": "👤", "content": _persona_content(pm.get_user_profile()), "editable": False},
        {"name": "social_anchors", "label": "社交锚点", "icon": "⚓", "content": _persona_content(pm.get_social_anchors()), "editable": False},
    ]
    return {"layers": layers, "empty": all(not l["content"] for l in layers)}


@app.get("/api/persona/layer/{name}")
async def api_persona_layer(name: str):
    hm = get_heimdall()
    pm = hm.provider.persona
    if not pm:
        raise HTTPException(404, "Persona manager not available")
    method_map = {
        "core_self": pm.get_core_self,
        "external_persona": pm.get_external_persona,
        "user_profile": pm.get_user_profile,
        "social_anchors": pm.get_social_anchors,
    }
    fn = method_map.get(name)
    if fn is None:
        raise HTTPException(404, f"Layer '{name}' not found")
    content = fn()
    return {"name": name, "content": _persona_content(content)}


@app.post("/api/persona/snapshot")
async def api_persona_snapshot():
    """Manually capture a persona snapshot."""
    hm = get_heimdall()
    pm = hm.provider.persona
    if not pm:
        raise HTTPException(404, "Persona manager not available")
    snap = _persona_snapshot.capture(pm)
    return {"snapshot_id": snap["id"], "timestamp": snap["timestamp"], "core_value_count": len(snap["core_self"])}


@app.get("/api/persona/drift")
async def api_persona_drift():
    """Return persona drift time-series data."""
    hm = get_heimdall()
    pm = hm.provider.persona
    if not pm:
        return {"snapshots": [], "current_drift": {}}
    return _persona_snapshot.get_drift_data(pm)


# ---------------------------------------------------------------------------
# Reconnect + Recent
# ---------------------------------------------------------------------------

@app.get("/api/reconnect")
async def api_reconnect():
    hm = get_heimdall()
    sg = hm.provider.social_graph
    if not sg:
        return {"suggestions": []}
    return {"suggestions": sg.get_reconnect_suggestions()}


# ---------------------------------------------------------------------------
# Social Perspective Retrieval (V3.0 section 5.2)
# ---------------------------------------------------------------------------

@app.get("/api/social/perspective")
async def api_social_perspective(
    entity_name: str = "",
    entity_id: str = "",
    days: int = 30,
):
    """Social-perspective retrieval with 3D trend analysis.

    Query social_graph materialized view → get relation metrics →
    retrieve associated memory_edges → calculate 3D trends →
    generate relationship insights.

    Example: /api/social/perspective?entity_name=配偶
    Returns: entity info, 3D metrics, trend analysis, insight narrative.
    """
    hm = get_heimdall()
    gs = hm.provider.store
    sg = hm.provider.social_graph

    if not gs or not gs._conn:
        raise HTTPException(503, "Entity store not available")

    # Resolve entity
    entity = None
    if entity_id:
        entity = gs.get_entity(entity_id)
    elif entity_name:
        entity = gs.get_entity_by_name(entity_name)
        if not entity:
            # Fuzzy search
            results = gs.search_entities(entity_name, limit=5)
            if results:
                entity = results[0]

    if not entity:
        raise HTTPException(404, f"Entity not found: {entity_name or entity_id}")

    eid = entity["entity_id"]
    now = time.time()
    cutoff_30d = now - (days * 86400)

    # 1. Get social graph connections for this entity
    connections = sg.get_connections(eid) if sg else []

    # 2. Get memory edges for interaction history
    memory_edges = []
    try:
        rows = gs._conn.execute(
            "SELECT me.*, e.display_name as entity_name, e.entity_type "
            "FROM heimdall_memory_edges me "
            "LEFT JOIN heimdall_entities e ON me.entity_id = e.entity_id "
            "WHERE me.entity_id = ? "
            "ORDER BY me.timestamp DESC LIMIT 100",
            (eid,),
        ).fetchall()
        memory_edges = [dict(r) for r in rows]
    except Exception:
        pass

    # 3. Calculate 3D social metrics
    recent_edges = [m for m in memory_edges if m.get("timestamp", 0) >= cutoff_30d]
    historical_edges = [m for m in memory_edges if m.get("timestamp", 0) < cutoff_30d]

    def _calc_metrics(edges: list) -> dict:
        if not edges:
            return {"count": 0, "avg_emotion": 0, "volatility": 0, "intensity": 0}
        emotions = [e.get("emotion", 0) or 0 for e in edges]
        avg_e = sum(emotions) / len(emotions)
        variance = sum((e - avg_e) ** 2 for e in emotions) / len(emotions) if emotions else 0
        return {
            "count": len(edges),
            "avg_emotion": round(avg_e, 4),
            "volatility": round(variance ** 0.5, 4),
            "intensity": round(1.0 / (1.0 + 2.71828 ** (-0.1 * len(edges))), 4),
        }

    recent_metrics = _calc_metrics(recent_edges)
    hist_metrics = _calc_metrics(historical_edges)

    # 4. Calculate trends
    def _trend(recent_val: float, hist_val: float) -> str:
        if hist_val == 0:
            return "new"
        delta = (recent_val - hist_val) / hist_val
        if delta > 0.15:
            return "rising"
        elif delta < -0.15:
            return "declining"
        return "stable"

    trends = {
        "intensity": {
            "recent": recent_metrics["intensity"],
            "historical": hist_metrics["intensity"],
            "trend": _trend(recent_metrics["intensity"], hist_metrics["intensity"]),
        },
        "valence": {
            "recent": recent_metrics["avg_emotion"],
            "historical": hist_metrics["avg_emotion"],
            "trend": _trend(recent_metrics["avg_emotion"], hist_metrics["avg_emotion"]),
        },
        "volatility": {
            "recent": recent_metrics["volatility"],
            "historical": hist_metrics["volatility"],
            "trend": _trend(recent_metrics["volatility"], hist_metrics["volatility"]),
        },
        "interaction_count": {
            "recent": recent_metrics["count"],
            "historical": hist_metrics["count"],
            "trend": _trend(recent_metrics["count"], hist_metrics["count"]),
        },
    }

    # 5. Generate insight narrative
    insight = _generate_social_insight(
        entity.get("display_name", "该联系人"),
        entity.get("entity_type", "person"),
        recent_metrics, hist_metrics, trends, days,
    )

    # 6. Reconnect check
    reconnect = None
    last_interaction = max(
        (e.get("timestamp", 0) for e in memory_edges), default=0
    )
    if last_interaction > 0:
        days_since = (now - last_interaction) / 86400
        if days_since > 90 and hist_metrics["intensity"] > 0.3:
            reconnect = {
                "days_inactive": round(days_since),
                "historical_intensity": hist_metrics["intensity"],
                "level": "L2",
                "suggestion": f"已经{round(days_since)}天未互动，但历史互动强度较高，建议主动联系",
            }

    # 7. Interaction timeline (last 20)
    timeline = []
    for me in memory_edges[:20]:
        ts = me.get("timestamp", 0)
        timeline.append({
            "timestamp": ts,
            "date": datetime.fromtimestamp(ts).strftime("%Y-%m-%d") if ts else "?",
            "role": me.get("role", ""),
            "emotion": me.get("emotion"),
        })

    return {
        "entity": {
            "entity_id": eid,
            "name": entity.get("display_name", "?"),
            "type": entity.get("entity_type", "person"),
            "occurrence_count": entity.get("occurrence_count", 0),
            "first_seen": entity.get("first_seen_at", 0),
            "last_seen": entity.get("last_seen_at", 0),
        },
        "connections": [
            {
                "target_id": c.get("target_entity_id") or c.get("entity_id", ""),
                "relationship": c.get("relationship_type", ""),
                "intensity": c.get("intensity", 0),
                "valence": c.get("valence", 0),
                "volatility": c.get("volatility", 0),
                "health_score": c.get("health_score", 0.5),
            }
            for c in connections[:10]
        ],
        "metrics": {
            "recent": recent_metrics,
            "historical": hist_metrics,
            "total_interactions": len(memory_edges),
        },
        "trends": trends,
        "insight": insight,
        "reconnect": reconnect,
        "timeline": timeline,
        "generated_at": datetime.now().isoformat(),
    }


def _generate_social_insight(
    name: str, etype: str,
    recent: dict, hist: dict, trends: dict, days: int,
) -> str:
    """Generate a human-readable social insight narrative.

    This produces template-based insights locally. For deep analysis,
    the /api/insights endpoint uses LLM generation.
    """
    parts = []

    # Intensity trend
    i_trend = trends["intensity"]["trend"]
    if i_trend == "declining":
        delta = 0
        if hist["intensity"] > 0:
            delta = (recent["intensity"] - hist["intensity"]) / hist["intensity"] * 100
        parts.append(f"最近{days}天互动频率下降了{abs(delta):.0f}%")
    elif i_trend == "rising":
        parts.append(f"最近{days}天互动频率上升")
    elif i_trend == "stable" and recent["count"] > 0:
        parts.append("互动频率保持稳定")

    # Valence trend
    v_trend = trends["valence"]["trend"]
    if v_trend == "declining":
        parts.append("情绪效价有所下降")
    elif v_trend == "rising":
        parts.append("单次互动质量上升，效价维持高位" if recent["avg_emotion"] > 0.3 else "情绪效价上升")
    elif v_trend == "stable" and recent["avg_emotion"] > 0.3:
        parts.append("情绪效价维持高位")

    # Volatility
    if recent["volatility"] > 0.4:
        parts.append("近期情绪波动较大")
    elif recent["volatility"] < 0.15 and recent["count"] > 0:
        parts.append("关系状态稳定")

    if not parts:
        if recent["count"] == 0:
            return f"最近{days}天没有与{name}的互动记录"
        return f"与{name}的关系保持正常"

    return f"最近{days}天：{'，'.join(parts)}"


# ---------------------------------------------------------------------------
# Resource Monitoring (V3.0 section 6.2)
# ---------------------------------------------------------------------------

@app.get("/api/system/resources")
async def api_system_resources():
    """Report current resource usage against dual-cabin budget.

    Hot cabin budget: ≤100MB (standard) / ≤125MB (high performance).
    Reports SQLite size, entity counts, degradation status.
    """
    import sys
    from hermes_constants import get_heimdall_home

    heimdall_dir = get_heimdall_home()
    db_path = heimdall_dir / "heimdall.db"

    # DB file sizes
    db_size = 0
    wal_size = 0
    shm_size = 0
    if db_path.exists():
        db_size = db_path.stat().st_size
    wal_path = heimdall_dir / "heimdall.db-wal"
    if wal_path.exists():
        wal_size = wal_path.stat().st_size
    shm_path = heimdall_dir / "heimdall.db-shm"
    if shm_path.exists():
        shm_size = shm_path.stat().st_size

    total_db_bytes = db_size + wal_size + shm_size

    # Entity and edge counts
    hm = get_heimdall()
    gs = hm.provider.store

    entity_count = gs.get_entity_count() if gs else 0
    memory_edge_count = 0
    social_edge_count = 0
    knowledge_count = 0
    try:
        if gs and gs._conn:
            memory_edge_count = gs._conn.execute(
                "SELECT COUNT(*) FROM heimdall_memory_edges"
            ).fetchone()[0]
            social_edge_count = gs._conn.execute(
                "SELECT COUNT(*) FROM heimdall_social_graph"
            ).fetchone()[0]
            knowledge_count = gs._conn.execute(
                "SELECT COUNT(*) FROM heimdall_knowledge_entries"
            ).fetchone()[0]
    except Exception:
        pass

    # Python process memory (approximate RSS)
    try:
        import psutil
        proc = psutil.Process()
        rss_mb = proc.memory_info().rss / (1024 * 1024)
    except ImportError:
        rss_mb = 0

    # Budget calculation
    budget_mb = 100
    high_perf = os.getenv("HEIMDALL_HIGH_PERF", "") == "1"
    if high_perf:
        budget_mb = 125

    usage_pct = (rss_mb / budget_mb * 100) if rss_mb > 0 else 0

    # Degradation status
    degradation = []
    remaining = budget_mb - rss_mb
    if remaining < 8:
        degradation.append("dynamic_embedding_cache_disabled")
    if remaining < 18:
        degradation.append("knowledge_index_trimmed")
    if remaining < 23:
        degradation.append("group_render_fallback_list_only")

    return {
        "budget": {
            "total_mb": budget_mb,
            "used_mb": round(rss_mb, 1),
            "remaining_mb": round(remaining, 1),
            "usage_pct": round(usage_pct, 1),
            "high_performance_mode": high_perf,
        },
        "storage": {
            "db_mb": round(db_size / (1024 * 1024), 2),
            "wal_mb": round(wal_size / (1024 * 1024), 2),
            "total_db_mb": round(total_db_bytes / (1024 * 1024), 2),
        },
        "counts": {
            "entities": entity_count,
            "memory_edges": memory_edge_count,
            "social_edges": social_edge_count,
            "knowledge_entries": knowledge_count,
        },
        "degradation": degradation,
        "python_version": sys.version.split()[0],
        "generated_at": datetime.now().isoformat(),
    }


# ---------------------------------------------------------------------------
# Pivot Moment Detection (V3.0 Phase 2)
# ---------------------------------------------------------------------------

@app.get("/api/insights/pivots")
async def api_pivots(limit: int = 10):
    """Detect pivot moments — bridge nodes with ≤7-day timestamp gaps.

    Returns scored pivot moments with narrative, domain change indicators,
    and emotion deltas. Top 3 by default for monthly summaries.
    """
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.views.pivots import PivotDetector

    detector = PivotDetector(gs)
    pivots = detector.detect(top_n=limit)

    return {
        "pivots": pivots,
        "total": len(pivots),
        "generated_at": datetime.now().isoformat(),
    }


# ---------------------------------------------------------------------------
# Memory Pollution Rollback (V3.0 Phase 2)
# ---------------------------------------------------------------------------

@app.get("/api/memories/flagged")
async def api_memories_flagged(limit: int = 50):
    """Get memory edges flagged for review (potential contamination)."""
    hm = get_heimdall()
    gs = hm.provider.store
    return {"flagged": gs.get_flagged_edges(limit=limit)}


@app.post("/api/memories/{edge_id}/flag")
async def api_memory_flag(edge_id: str, request: Request):
    """Flag a memory edge for review."""
    body = await request.json()
    reason = body.get("reason", "")
    hm = get_heimdall()
    gs = hm.provider.store
    ok = gs.flag_memory_edge(edge_id, reason=reason)
    return {"flagged": ok, "edge_id": edge_id}


@app.post("/api/memories/{edge_id}/unflag")
async def api_memory_unflag(edge_id: str):
    """Remove flag from a memory edge."""
    hm = get_heimdall()
    gs = hm.provider.store
    ok = gs.unflag_memory_edge(edge_id)
    return {"unflagged": ok, "edge_id": edge_id}


@app.post("/api/entities/{entity_id}/cascade-delete")
async def api_cascade_delete(entity_id: str):
    """Cascade delete an entity: audit → delete edges → mark dormant.

    V3.0 Phase 2: preserves audit trail, marks entity dormant rather
    than hard-deleting it.
    """
    hm = get_heimdall()
    gs = hm.provider.store

    entity = gs.get_entity(entity_id)
    if not entity:
        raise HTTPException(404, "Entity not found")

    summary = gs.cascade_delete_entity(entity_id)
    return {
        "entity_id": entity_id,
        "entity_name": entity.get("display_name", "?"),
        "summary": summary,
        "timestamp": datetime.now().isoformat(),
    }


@app.post("/api/entities/{entity_id}/reactivate")
async def api_reactivate_entity(entity_id: str):
    """Reactivate a dormant entity."""
    hm = get_heimdall()
    gs = hm.provider.store

    entity = gs.get_entity(entity_id)
    if not entity:
        raise HTTPException(404, "Entity not found")
    if entity.get("status") != "dormant":
        raise HTTPException(400, "Entity is not dormant")

    ok = gs.reactivate_entity(entity_id)
    return {"reactivated": ok, "entity_id": entity_id}


@app.get("/api/entities/dormant")
async def api_dormant_entities(limit: int = 50):
    """List dormant entities (marked for cleanup)."""
    hm = get_heimdall()
    gs = hm.provider.store
    return {"entities": gs.get_dormant_entities(limit=limit)}


@app.get("/api/memories/recent")
async def api_memories_recent(limit: int = 20, session_id: str = "", search: str = ""):
    hm = get_heimdall()
    gs = hm.provider.store
    sid = ""
    if session_id == "current":
        sid = hm.provider._session_id if hasattr(hm.provider, "_session_id") else ""
    edges = []
    try:
        where = []
        params = []
        if sid:
            where.append("me.session_id = ?")
            params.append(sid)
        if search:
            where.append("(e.display_name LIKE ? OR me.role LIKE ?)")
            params.extend([f"%{search}%", f"%{search}%"])
        where_clause = f"WHERE {' AND '.join(where)}" if where else ""
        cur = gs._conn.execute(
            f"SELECT me.id, me.entity_id, me.role, me.emotion, me.timestamp, me.session_id "
            f"FROM heimdall_memory_edges me "
            f"LEFT JOIN heimdall_entities e ON me.entity_id = e.entity_id "
            f"{where_clause} "
            f"ORDER BY me.timestamp DESC LIMIT ?",
            params + [limit],
        )
        for row in cur:
            edges.append({
                "id": row[0],
                "entity_id": row[1],
                "role": row[2],
                "emotion": row[3],
                "timestamp": row[4],
                "session_id": row[5],
            })
        # Resolve entity names in a second pass
        eids = [e["entity_id"] for e in edges if e["entity_id"]]
        if eids:
            entities = {}
            placeholders = ",".join("?" for _ in eids)
            try:
                erows = gs._conn.execute(
                    f"SELECT entity_id, display_name FROM heimdall_entities WHERE entity_id IN ({placeholders})",
                    eids,
                ).fetchall()
                for r in erows:
                    entities[r["entity_id"]] = r["display_name"]
            except Exception:
                pass
            for e in edges:
                e["entity_name"] = entities.get(e["entity_id"], "?")
    except Exception:
        pass
    return {"edges": edges}


@app.get("/api/memories/session")
async def api_memories_session(limit: int = 50):
    """Return memory edges for the current session only (short-term)."""
    hm = get_heimdall()
    gs = hm.provider.store
    sid = hm.provider._session_id if hasattr(hm.provider, "_session_id") else ""
    edges = []
    if gs._conn and sid:
        try:
            cur = gs._conn.execute(
                "SELECT me.id, me.entity_id, me.role, me.emotion, me.timestamp, e.display_name, e.entity_type "
                "FROM heimdall_memory_edges me "
                "LEFT JOIN heimdall_entities e ON me.entity_id = e.entity_id "
                "WHERE me.session_id = ? "
                "ORDER BY me.timestamp DESC LIMIT ?",
                (sid, limit),
            )
            for row in cur:
                edges.append({
                    "id": row[0], "entity_id": row[1], "role": row[2],
                    "emotion": row[3], "timestamp": row[4],
                    "entity_name": row[5] or "?", "entity_type": row[6] or "concept",
                })
        except Exception:
            pass
    return {"session_id": sid, "edges": edges}


@app.get("/api/memories/long-term")
async def api_memories_long_term(
    search: str = "",
    entity_type: str = "",
    emotion: str = "",
    limit: int = 30,
    offset: int = 0,
):
    """Browse long-term memories (excluding current session) with filters."""
    hm = get_heimdall()
    gs = hm.provider.store
    sid = hm.provider._session_id if hasattr(hm.provider, "_session_id") else ""

    if not gs._conn:
        return {"total": 0, "limit": limit, "offset": offset, "edges": []}

    where = ["me.session_id != ?"]
    params = [sid]

    if entity_type:
        where.append("e.entity_type = ?")
        params.append(entity_type)
    if emotion == "positive":
        where.append("me.emotion > 0.3")
    elif emotion == "negative":
        where.append("me.emotion < -0.3")
    elif emotion == "neutral":
        where.append("me.emotion BETWEEN -0.3 AND 0.3")
    if search:
        where.append("(e.display_name LIKE ? OR me.role LIKE ?)")
        params.extend([f"%{search}%", f"%{search}%"])

    where_clause = " AND ".join(where)

    try:
        count_row = gs._conn.execute(
            f"SELECT COUNT(*) as cnt FROM heimdall_memory_edges me "
            f"LEFT JOIN heimdall_entities e ON me.entity_id = e.entity_id "
            f"WHERE {where_clause}", params,
        ).fetchone()
        total = count_row["cnt"] if count_row else 0

        cur = gs._conn.execute(
            f"SELECT me.id, me.entity_id, me.role, me.emotion, me.timestamp, me.session_id, "
            f"e.display_name, e.entity_type "
            f"FROM heimdall_memory_edges me "
            f"LEFT JOIN heimdall_entities e ON me.entity_id = e.entity_id "
            f"WHERE {where_clause} "
            f"ORDER BY me.timestamp DESC LIMIT ? OFFSET ?",
            params + [limit, offset],
        )
        edges = []
        for row in cur:
            edges.append({
                "id": row[0], "entity_id": row[1], "role": row[2],
                "emotion": row[3], "timestamp": row[4], "session_id": row[5],
                "entity_name": row[6] or "?", "entity_type": row[7] or "concept",
            })
        return {"total": total, "limit": limit, "offset": offset, "edges": edges}
    except Exception as e:
        logger.warning("Long-term memory query failed: %s", e)
        return {"total": 0, "limit": limit, "offset": offset, "edges": []}


# ---------------------------------------------------------------------------
# Entities + Social Graph
# ---------------------------------------------------------------------------

@app.get("/api/entities")
async def api_entities(
    entity_type: str = "",
    track: str = "",
    session_id: str = "",
    limit: int = 50,
    offset: int = 0,
):
    hm = get_heimdall()
    gs = hm.provider.store
    if session_id == "current":
        session_id = hm.provider._session_id if hasattr(hm.provider, "_session_id") else ""
    entities = gs.list_entities(
        entity_type=entity_type or None,
        source_track=track or None,
        limit=limit,
        offset=offset,
    )
    if session_id:
        entities = [e for e in entities if e.get("source_session_id") == session_id]
    return {"entities": entities, "total": gs.get_entity_count()}


@app.get("/api/entities/search")
async def api_entities_search(q: str = "", limit: int = 10):
    hm = get_heimdall()
    gs = hm.provider.store
    return {"entities": gs.search_entities(q, limit=limit)}


@app.get("/api/entities/{entity_id}")
async def api_entity_detail(entity_id: str):
    hm = get_heimdall()
    gs = hm.provider.store
    sg = hm.provider.social_graph
    entity = gs.get_entity(entity_id)
    if not entity:
        raise HTTPException(404, "Entity not found")
    connections = sg.get_connections(entity_id) if sg else []
    return {"entity": entity, "connections": connections}


# Compatibility alias — desktop client uses /api/knowledge/entity/{id}
@app.get("/api/knowledge/entity/{entity_id}")
async def api_knowledge_entity_detail(entity_id: str):
    """Entity detail — compatibility route for desktop knowledge panel."""
    return await api_entity_detail(entity_id)


@app.get("/api/entities/{entity_id}/graph")
async def api_entity_graph(entity_id: str):
    """Return vis.js graph data for the entity's ego network (true 2-hop).

    Includes entity nodes, knowledge entry nodes, and three edge types:
    - social: from social_graph connections (solid line, valence-colored)
    - knowledge: from knowledge_edges (dashed line, blue)
    - memory: from memory_edges (dotted line, valence-colored)
    """
    hm = get_heimdall()
    gs = hm.provider.store
    sg = hm.provider.social_graph

    entity = gs.get_entity(entity_id)
    if not entity:
        raise HTTPException(404, "Entity not found")

    # Load community labels for annotation
    community_labels: dict = {}
    try:
        from heimdall.core.community import CommunityDetector
        detector = CommunityDetector(gs)
        communities = detector.get_communities()
        for c in communities:
            cid = c["community_id"]
            label = f"{c['emoji']} {c['label']}"
            for ent in c.get("entities", []):
                community_labels[ent["entity_id"]] = label
    except Exception:
        pass

    nodes = {entity_id: _entity_to_vis_node(entity)}
    edges = []
    _edge_ids = set()  # dedup edges

    def _with_community_label(ent: dict) -> dict:
        """Inject community_label from pre-loaded map."""
        eid = ent.get("entity_id", "")
        label = community_labels.get(eid, "")
        if label:
            ent = dict(ent)
            ent["community_label"] = label
        return ent

    def _add_edge(src, tgt, label, etype, intensity=0.5, valence=0):
        eid = f"{src}|{tgt}|{etype}"
        if eid in _edge_ids:
            return
        _edge_ids.add(eid)
        edge = {
            "from": src, "to": tgt, "label": label,
            "value": intensity,
            "edge_type": etype,
        }
        if etype == "knowledge":
            edge["dashes"] = True
            edge["color"] = {"color": "#3b82f6", "highlight": "#60a5fa"}
        elif etype == "social":
            edge["dashes"] = False
            edge["color"] = _valence_color(valence)
        else:
            edge["dashes"] = [4, 4]
            edge["color"] = _valence_color(valence)
        edges.append(edge)

    # --- 1-hop: social graph connections ---
    hop1_ids = set()
    connections = sg.get_connections(entity_id) if sg else []
    for conn in connections:
        # Social edges have source_entity_id / target_entity_id — pick the
        # one that is NOT the center entity to get the neighbor.
        src_id = conn.get("source_entity_id", "")
        tgt_id = conn.get("target_entity_id", "")
        tid = src_id if tgt_id == entity_id else tgt_id
        if not tid or tid == entity_id:
            continue
        target = gs.get_entity(tid)
        if target and tid not in nodes:
            nodes[tid] = _entity_to_vis_node(_with_community_label(target))
        if tid not in nodes:
            nodes[tid] = {"id": tid, "label": tid[:8], "group": "unknown"}
        _add_edge(entity_id, tid, "", "social",
                  conn.get("intensity", 0.5), conn.get("valence", 0))
        hop1_ids.add(tid)

    # --- Knowledge edges for center entity ---
    if gs._conn:
        try:
            krows = gs._conn.execute(
                "SELECT ke.entry_id, ke.title, ke.domain, ke.mastery_level "
                "FROM heimdall_knowledge_edges kee "
                "JOIN heimdall_knowledge_entries ke ON ke.entry_id = kee.entry_id "
                "WHERE kee.entity_id = ? LIMIT 30",
                (entity_id,),
            ).fetchall()
        except Exception:
            krows = []
        for kr in krows:
            kid = f"k_{kr['entry_id']}"
            if kid not in nodes:
                nodes[kid] = {
                    "id": kid, "label": (kr["title"] or kr["entry_id"])[:20],
                    "group": "knowledge", "color": "#3b82f6",
                    "shape": "square",
                    "title": f"Knowledge: {kr['title']}\ndomain: {kr.get('domain','')}\nmastery: {kr.get('mastery_level','')}",
                }
            _add_edge(entity_id, kid, kr.get("domain", "knowledge"), "knowledge", 0.6, 0)

    # --- 2-hop: connections of connected entities ---
    for tid in list(hop1_ids)[:8]:
        conns2 = sg.get_connections(tid) if sg else []
        for conn in conns2[:5]:
            s2 = conn.get("source_entity_id", "")
            t2 = conn.get("target_entity_id", "")
            nid = s2 if t2 == tid else t2
            if not nid or nid == tid or nid == entity_id:
                continue
            if nid in nodes:
                _add_edge(tid, nid, "", "social",
                          conn.get("intensity", 0.3), conn.get("valence", 0))
            elif len(nodes) < 80:
                target2 = gs.get_entity(nid)
                if target2:
                    nodes[nid] = _entity_to_vis_node(_with_community_label(target2))
                else:
                    nodes[nid] = {"id": nid, "label": nid[:8], "group": "unknown"}
                _add_edge(tid, nid, "", "social",
                          conn.get("intensity", 0.3), conn.get("valence", 0))

    # --- Memory edges for entities in graph ---
    if gs._conn:
        eids = list(nodes.keys())[:30]
        placeholders = ",".join("?" for _ in eids)
        try:
            mrows = gs._conn.execute(
                f"SELECT entity_id, role, emotion FROM heimdall_memory_edges "
                f"WHERE entity_id IN ({placeholders}) LIMIT 100",
                eids,
            ).fetchall()
        except Exception:
            mrows = []
        for mr in mrows:
            mr = dict(mr)
            eid = mr["entity_id"]
            if eid in nodes:
                nodes[eid]["title"] = (nodes[eid].get("title", "") or "") + f"\nrole: {mr['role']} emotion: {mr['emotion']:.2f}" if mr.get("emotion") is not None else f"\nrole: {mr['role']}"

    return {"nodes": list(nodes.values()), "edges": edges}


def _entity_to_vis_node(entity: dict) -> dict:
    # Artistic muted palette — warm/cool balanced, visually harmonious
    color_map = {
        "person":       {"background": "#e07a5f", "border": "#c0624a"},
        "organization": {"background": "#5e8c7e", "border": "#4a7265"},
        "project":      {"background": "#667bc0", "border": "#5065a8"},
        "tool":         {"background": "#7d9f6e", "border": "#658556"},
        "concept":      {"background": "#d4a76a", "border": "#bc8f52"},
        "skill":        {"background": "#b08abb", "border": "#9670a3"},
        "event":        {"background": "#6dacb5", "border": "#55909a"},
        "location":     {"background": "#d4956b", "border": "#bc7d53"},
        "media":        {"background": "#a089b5", "border": "#86709d"},
    }
    fallback = {"background": "#a09988", "border": "#888070"}
    etype = entity.get("entity_type", "concept")
    node = {
        "id": entity.get("entity_id", ""),
        "label": entity.get("display_name", "?")[:20],
        "group": etype,
        "color": color_map.get(etype, fallback),
        "title": f"{entity.get('display_name','?')}\ntype: {etype}\nseen: {entity.get('occurrence_count',0)}x\nconfidence: {entity.get('confidence',0):.0%}",
    }
    # Bridge node: diamond shape + border
    if entity.get("is_bridge"):
        node["shape"] = "diamond"
        node["borderWidth"] = 3
        base = color_map.get(etype, fallback)
        node["color"] = {"background": base["background"], "border": "#ec4899"}
    # Community coloring — assign stable color from community_id
    cid = entity.get("community_id")
    if cid is not None and not entity.get("is_bridge"):
        # 12-color palette designed for visual distinguishability
        community_palette = [
            "#667bc0", "#b08abb", "#e07a5f", "#d4956b", "#5e8c7e",
            "#d4a76a", "#6dacb5", "#7d9f6e", "#c47e8b", "#6b9eac",
            "#c0a060", "#9b7eb8",
        ]
        bg = community_palette[cid % len(community_palette)]
        node["color"] = {"background": bg, "border": bg}
        # Add community label to tooltip
        c_label = entity.get("community_label", "")
        if c_label:
            node["title"] = (node.get("title", "") or "") + f"\ncommunity: {c_label}"
    return node


def _valence_color(valence: float) -> dict:
    if valence > 0.3:
        return {"color": "#2ECC71", "highlight": "#27AE60"}
    elif valence < -0.3:
        return {"color": "#E74C3C", "highlight": "#C0392B"}
    else:
        return {"color": "#BDC3C7", "highlight": "#95A5A6"}


# ---------------------------------------------------------------------------
# Knowledge / Skills (学习的)
# ---------------------------------------------------------------------------

@app.get("/api/knowledge")
async def api_knowledge(q: str = "", domain: str = "", limit: int = 20):
    hm = get_heimdall()
    km = hm.provider.knowledge
    if not km:
        return {"entries": []}
    if q:
        entries = km.search(q, domain=domain or None, limit=limit)
    else:
        entries = km.search("", limit=limit)
    return {"entries": entries}


@app.get("/api/knowledge/stale")
async def api_knowledge_stale(days: int = 30):
    hm = get_heimdall()
    km = hm.provider.knowledge
    if not km:
        return {"entries": []}
    return {"entries": km.get_stale_entries(days)}


@app.get("/api/knowledge/search")
async def api_knowledge_search(query: str = "", limit: int = 10):
    """Search entities across all namespaces — desktop IPC alias.

    Returns results formatted with field names the desktop expects
    (id/name/type/description/importance).
    """
    hm = get_heimdall()
    gs = hm.provider.store
    if not gs:
        return []
    raw = gs.search_entities(query, limit=limit)
    return [
        {
            "id": e.get("entity_id", e.get("id", "")),
            "name": e.get("name", ""),
            "type": (e.get("types") or ["concept"])[0],
            "description": e.get("description", ""),
            "importance": e.get("importance", 0.5),
        }
        for e in raw
    ]


@app.get("/api/knowledge/overview")
async def api_knowledge_overview(namespace: str = None):
    """Return knowledge overview text for a namespace.

    Reads from OverviewCache via DynamicContextBuilder.
    """
    hm = get_heimdall()
    builder = getattr(hm.provider, "_dynamic_builder", None)
    store = hm.provider.store
    if not builder or not store:
        return {"overview": "", "entity_count": 0}
    text = builder.build(namespace)
    count = store.count_by_namespace(namespace) if namespace else store.count_all()
    return {"overview": text, "entity_count": count}


@app.get("/api/skills")
async def api_skills():
    hm = get_heimdall()
    km = hm.provider.knowledge
    if not km:
        return {"skills": []}
    return {"skills": km.get_skills_needing_review()}


@app.post("/api/skills/{skill_id}/mastery")
async def api_update_mastery(skill_id: str, request: Request):
    body = await request.json()
    level = body.get("mastery_level", "了解")
    hm = get_heimdall()
    km = hm.provider.knowledge
    if not km:
        raise HTTPException(404)
    result = km.update_mastery(skill_id, level, body.get("parent_domain", ""))
    return {"result": result}


@app.post("/api/skills/suggest")
async def api_skills_suggest(request: Request):
    """AI-assisted skill mastery upgrade suggestions."""
    body = await request.json()
    skill_name = body.get("skill_name", "")
    hm = get_heimdall()
    gs = hm.provider.store
    if not gs or not gs._conn:
        return {"suggestions": [], "error": "Store not available"}

    # Fetch all skill mastery records
    rows = gs._conn.execute(
        "SELECT skill_name, parent_domain, mastery_level FROM heimdall_skill_mastery ORDER BY last_interacted DESC LIMIT 30"
    ).fetchall()
    skills = [dict(r) for r in rows]
    if skill_name:
        skills = [s for s in skills if s.get("skill_name") == skill_name]
    if not skills:
        return {"suggestions": [], "message": "No skills to analyze"}

    skills_text = "\n".join(
        f"- {s['skill_name']}: mastery={s.get('mastery_level','了解')}, domain={s.get('parent_domain','')}"
        for s in skills[:15]
    )
    prompt = (
        "分析以下技能列表，为每个技能建议掌握度升级。掌握度等级: 不了解 → 了解 → 练习中 → 掌握 → 精通\n\n"
        f"当前技能:\n{skills_text}\n\n"
        "返回 JSON 数组（仅 JSON，不要其他文字）:\n"
        '[{"skill_name": "...", "current_level": "...", "suggested_level": "...", "reason": "...", "confidence": 0.0}]\n\n'
        "规则:\n- 如果当前已经是'精通'，不建议升级\n- confidence 是 0-1 的建议置信度\n- reason 简短（20字内）\n- 最多建议3个技能"
    )

    try:
        provider, model = _get_model_config()
        from agent.auxiliary_client import call_llm
        resp = call_llm(
            provider=provider, model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3, max_tokens=800, timeout=30,
        )
        text = resp.choices[0].message.content
        match = re.search(r'\[.*\]', text, re.DOTALL)
        suggestions = json.loads(match.group(0)) if match else []
        return {"suggestions": suggestions, "generated_at": datetime.now().isoformat()}
    except Exception as e:
        logger.warning("Skill suggest failed: %s", e)
        return {"suggestions": [], "error": str(e)}


# ---------------------------------------------------------------------------
# Community Detection (分组) — V3.0 section 3.1
# ---------------------------------------------------------------------------

@app.post("/api/communities/detect")
async def api_communities_detect():
    """Trigger GMM-based community detection on the entity graph.

    Only activates when cold-start thresholds are met:
    - ≥50 active entities
    - ≥5 entities per potential community
    """
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.core.community import CommunityDetector

    detector = CommunityDetector(gs)
    assignments = detector.detect()
    communities = detector.get_communities()
    bridges = detector.get_bridge_entities()

    classified = sum(1 for v in assignments.values() if v >= 0)
    unclassified = len(assignments) - classified

    # Persist community metadata to file cache so overview rebuild can use it
    metadata = detector.get_communities_metadata()
    if metadata and hm.provider._community_cache:
        for ns_info in gs.list_namespaces():
            ns = ns_info.get("namespace", "general")
            try:
                hm.provider._community_cache.write(ns, metadata)
            except Exception:
                pass

    return {
        "communities": communities,
        "bridges": [dict(b) for b in bridges],
        "stats": {
            "total_entities": len(assignments),
            "classified": classified,
            "unclassified": unclassified,
            "community_count": len(communities),
            "bridge_count": len(bridges),
        },
        "generated_at": datetime.now().isoformat(),
    }


@app.get("/api/communities")
async def api_communities():
    """List discovered communities with their entity members."""
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.core.community import CommunityDetector

    detector = CommunityDetector(gs)
    communities = detector.get_communities()
    bridges = detector.get_bridge_entities()

    return {
        "communities": communities,
        "bridges": [dict(b) for b in bridges],
        "cold_start_met": gs.get_entity_count() >= 50,
    }


# ---------------------------------------------------------------------------
# PageRank (V3.0 section 3.2)
# ---------------------------------------------------------------------------

@app.post("/api/pagerank/recalc")
async def api_pagerank_recalc():
    """Recalculate PageRank for all active entities (2-hop propagation).

    Applies base_type_weight priors, forgetting curve, and error monitoring.
    Returns list of entities with elevated pagerank_error needing cloud sync.
    """
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.core.pagerank import PageRankComputer

    pr = PageRankComputer(gs)
    results = pr.recalc_all()
    needing_sync = pr.get_entities_needing_sync()

    return {
        "updated": len(results),
        "score_range": {
            "min": round(min(results.values()), 4) if results else 0,
            "max": round(max(results.values()), 4) if results else 0,
        },
        "needing_sync": len(needing_sync),
        "sync_candidates": [
            {"entity_id": e["entity_id"], "name": e["display_name"],
             "pagerank": e["pagerank"], "error": e["pagerank_error"]}
            for e in needing_sync[:10]
        ],
        "generated_at": datetime.now().isoformat(),
    }


@app.get("/api/pagerank/sync-candidates")
async def api_pagerank_sync_candidates():
    """List entities with PageRank error exceeding 0.2 threshold."""
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.core.pagerank import PageRankComputer

    pr = PageRankComputer(gs)
    return {"entities": [
        {"entity_id": e["entity_id"], "name": e["display_name"],
         "pagerank": e["pagerank"], "error": e["pagerank_error"]}
        for e in pr.get_entities_needing_sync()
    ]}


# ---------------------------------------------------------------------------
# Importance Scoring (P2.3)
# ---------------------------------------------------------------------------

@app.post("/api/importance/recalc")
async def api_importance_recalc(namespace: Optional[str] = None):
    """Recalculate composite importance scores for all active entities.

    Formula: 0.30×confidence + 0.30×pagerank + 0.20×occurrence
           + 0.15×recency + 0.05×bridge_score

    Returns count of updated entities and level distribution.
    """
    hm = get_heimdall()
    store = hm.provider.store
    updated = store.recalculate_all_importance(namespace=namespace)
    levels = store.get_importance_level_counts(namespace=namespace)
    return {
        "updated": updated,
        "levels": levels,
        "generated_at": datetime.now().isoformat(),
    }


@app.get("/api/importance/top")
async def api_importance_top(
    limit: int = 20,
    namespace: Optional[str] = None,
):
    """Return top-N entities ranked by composite importance score."""
    hm = get_heimdall()
    store = hm.provider.store
    entities = store.get_high_importance_entities(limit=limit, namespace=namespace)
    return {
        "entities": entities,
        "generated_at": datetime.now().isoformat(),
    }


@app.get("/api/importance/levels")
async def api_importance_levels(namespace: Optional[str] = None):
    """Return entity count per importance level (critical/high/medium/low)."""
    hm = get_heimdall()
    store = hm.provider.store
    return {
        "levels": store.get_importance_level_counts(namespace=namespace),
        "generated_at": datetime.now().isoformat(),
    }


# ---------------------------------------------------------------------------
# Summary Tree (P2.1)
# ---------------------------------------------------------------------------

@app.post("/api/summary/bootstrap")
async def api_summary_bootstrap(
    target_date: Optional[str] = None,
    namespace: str = "general",
):
    """Create a daily (L0) summary entry for the given date (default: today)."""
    hm = get_heimdall()
    store = hm.provider.store
    store.summary_tree_init_schema()
    from datetime import date as dt
    d = dt.fromisoformat(target_date) if target_date else None
    sid = store.summary_tree_bootstrap_daily(target_date=d, namespace=namespace)
    return {
        "created": sid is not None,
        "summary_id": sid,
        "target_date": target_date or dt.today().isoformat(),
        "generated_at": datetime.now().isoformat(),
    }


@app.post("/api/summary/cascade")
async def api_summary_cascade(namespace: str = "general"):
    """Run cascade check: promote summaries from lower to higher levels."""
    hm = get_heimdall()
    store = hm.provider.store
    store.summary_tree_init_schema()
    results = store.summary_tree_cascade(namespace=namespace)
    return {
        "cascade_results": results,
        "generated_at": datetime.now().isoformat(),
    }


@app.get("/api/summary/list")
async def api_summary_list(
    level: Optional[int] = None,
    namespace: str = "general",
    limit: int = 50,
):
    """List summary entries, optionally filtered by level."""
    hm = get_heimdall()
    store = hm.provider.store
    store.summary_tree_init_schema()
    summaries = store.get_summaries(level=level, namespace=namespace, limit=limit)
    return {
        "summaries": summaries,
        "count": len(summaries),
        "generated_at": datetime.now().isoformat(),
    }


@app.get("/api/summary/latest")
async def api_summary_latest(level: int = 1, namespace: str = "general"):
    """Get the latest completed summary at the given level."""
    hm = get_heimdall()
    store = hm.provider.store
    store.summary_tree_init_schema()
    summary = store.get_latest_summary(level=level, namespace=namespace)
    return {
        "summary": summary,
        "generated_at": datetime.now().isoformat(),
    }


@app.post("/api/summary/generate")
async def api_summary_generate(summary_id: int):
    """Generate LLM summary for a pending summary_tree entry."""
    hm = get_heimdall()
    store = hm.provider.store
    store.summary_tree_init_schema()
    ok = store.summary_tree_generate(summary_id)
    return {"generated": ok, "summary_id": summary_id}


# ---------------------------------------------------------------------------
# Serendipity Search (V3.0 section 3.3)
# ---------------------------------------------------------------------------

@app.get("/api/search/serendipity")
async def api_search_serendipity(
    q: str = "",
    entity_id: str = "",
    limit: int = 10,
):
    """Retrieval with serendipity-weighted ranking (phase 1 fixed weights).

    Combines vector_sim(0.5) + pagerank(0.2) + time_decay(0.2) +
    freshness(0.1) + serendipity(0.15) for ranked results.
    """
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.core.pagerank import PageRankComputer

    pr = PageRankComputer(gs)

    # Start with entity search or recent entities
    if q:
        candidates = gs.search_entities(q, limit=limit * 2)
        if not candidates:
            # Fallback: scan by display_name LIKE
            try:
                rows = gs._conn.execute(
                    "SELECT * FROM heimdall_entities WHERE display_name LIKE ? "
                    "AND status = 'active' ORDER BY last_seen_at DESC LIMIT ?",
                    (f"%{q}%", limit * 2),
                ).fetchall()
                candidates = [dict(r) for r in rows]
            except Exception:
                candidates = gs.list_entities(limit=limit * 2)
    elif entity_id:
        entity = gs.get_entity(entity_id)
        candidates = [entity] if entity else []
    else:
        candidates = gs.list_entities(limit=limit * 2)

    # Score each candidate with retrieval weights
    scored = []
    for c in candidates:
        eid = c["entity_id"]
        score = pr.get_retrieval_score(eid, vector_sim=0.5 if q else 0.3)
        scored.append({
            "entity_id": eid,
            "name": c.get("display_name", "?"),
            "entity_type": c.get("entity_type", "concept"),
            "pagerank": round(c.get("pagerank", 1.0), 4),
            "is_bridge": bool(c.get("is_bridge")),
            "bridge_score": round(c.get("bridge_score", 0.0), 4),
            "community_id": c.get("community_id"),
            "retrieval_score": round(score, 4),
            "weights_used": {
                k: round(v, 4)
                for k, v in {"vector_sim": 0.5, "pagerank": 0.2,
                             "time_decay": 0.2, "freshness": 0.1,
                             "serendipity": 0.15}.items()
            },
        })

    scored.sort(key=lambda x: x["retrieval_score"], reverse=True)
    return {"results": scored[:limit], "query": q, "total_candidates": len(candidates)}


# ---------------------------------------------------------------------------
# FTRL Online Learning (V3.0 section 3.4)
# ---------------------------------------------------------------------------

@app.post("/api/ftrl/feedback")
async def api_ftrl_feedback(request: Request):
    """Record user feedback on a retrieval result for FTRL online learning.

    Body: {entity_id, features: [6 floats], clicked: bool}
    Phase 1: accumulates feedback, transitions to FTRL after ≥20 samples.
    """
    body = await request.json()
    entity_id = body.get("entity_id", "")
    features = body.get("features", [])
    clicked = body.get("clicked", True)

    if not features or len(features) < 6:
        # Auto-build features from entity
        hm = get_heimdall()
        gs = hm.provider.store
        entity = gs.get_entity(entity_id) if entity_id else None
        if entity:
            features = [
                body.get("vector_sim", 0.5),
                min(entity.get("pagerank", 1.0) / 10.0, 1.0),
                body.get("time_decay", 0.5),
                entity.get("bridge_score", 0.0),
                body.get("freshness", 0.5),
                body.get("serendipity", 0.15),
            ]
        else:
            raise HTTPException(400, "No entity_id or features provided")

    from heimdall.core.ftrl import get_ftrl
    from hermes_constants import get_heimdall_home

    model_path = get_heimdall_home() / "ftrl_state.json"
    ftrl = get_ftrl(model_path)
    ftrl.record_feedback(features, clicked=clicked)
    ftrl.save()

    return {
        "feedback_count": ftrl._feedback_count,
        "phase1": ftrl._phase1,
        "current_weights": ftrl.get_weights(),
        "recorded_at": datetime.now().isoformat(),
    }


@app.get("/api/ftrl/state")
async def api_ftrl_state():
    """Return current FTRL model state and weights."""
    from heimdall.core.ftrl import get_ftrl
    from hermes_constants import get_heimdall_home

    model_path = get_heimdall_home() / "ftrl_state.json"
    ftrl = get_ftrl(model_path)

    return {
        "phase": "phase1 (fixed weights)" if ftrl._phase1 else "phase2 (FTRL learned)",
        "dim": ftrl.dim,
        "feature_names": ftrl.feature_names,
        "feedback_count": ftrl._feedback_count,
        "weights": ftrl.get_weights(),
        "alpha": ftrl.alpha,
        "beta": ftrl.beta,
        "l1": ftrl.l1,
        "l2": ftrl.l2,
    }


@app.get("/api/communities/bridges")
async def api_communities_bridges():
    """List bridge entities connecting multiple communities."""
    hm = get_heimdall()
    gs = hm.provider.store
    from heimdall.core.community import CommunityDetector

    detector = CommunityDetector(gs)
    return {"bridges": [dict(b) for b in detector.get_bridge_entities()]}


# ---------------------------------------------------------------------------
# Summary (总结)
# ---------------------------------------------------------------------------

@app.get("/api/summary")
async def api_summary():
    hm = get_heimdall()
    gs = hm.provider.store
    sg = hm.provider.social_graph
    km = hm.provider.knowledge

    type_counts = {}
    try:
        cur = gs._conn.execute("SELECT entity_type, COUNT(*) FROM heimdall_entities WHERE status='active' GROUP BY entity_type")
        for row in cur:
            type_counts[row[0]] = row[1]
    except Exception:
        pass

    return {
        "entity_total": gs.get_entity_count(),
        "entity_types": type_counts,
        "social_stats": sg.get_stats() if sg else {},
        "knowledge_count": len(km.search("", limit=1000)) if km else 0,
        "reconnect_count": len(sg.get_reconnect_suggestions()) if sg else 0,
    }


@app.get("/api/insights")
async def api_insights():
    """Generate micro-AGI narrative insights from current state."""
    hm = get_heimdall()
    gs = hm.provider.store
    sg = hm.provider.social_graph
    km = hm.provider.knowledge

    # Aggregate data snapshot
    entity_total = gs.get_entity_count()
    type_counts = {}
    try:
        cur = gs._conn.execute(
            "SELECT entity_type, COUNT(*) FROM heimdall_entities WHERE status='active' GROUP BY entity_type"
        )
        for row in cur:
            type_counts[row[0]] = row[1]
    except Exception:
        pass

    # Top entities by occurrence
    top_entities = []
    try:
        cur = gs._conn.execute(
            "SELECT display_name, entity_type, occurrence_count FROM heimdall_entities "
            "WHERE status='active' ORDER BY occurrence_count DESC LIMIT 5"
        )
        top_entities = [dict(r) for r in cur.fetchall()]
    except Exception:
        pass

    social_stats = sg.get_stats() if sg else {}
    knowledge_count = len(km.search("", limit=1000)) if km else 0
    reconnect_count = len(sg.get_reconnect_suggestions()) if sg else 0

    # Recent memories count
    memory_count = 0
    try:
        row = gs._conn.execute("SELECT COUNT(*) as cnt FROM heimdall_memory_edges").fetchone()
        if row:
            memory_count = row["cnt"]
    except Exception:
        pass

    data_snapshot = {
        "entity_total": entity_total,
        "entity_types": type_counts,
        "top_entities": top_entities,
        "social_stats": social_stats,
        "knowledge_count": knowledge_count,
        "reconnect_count": reconnect_count,
        "memory_count": memory_count,
    }

    prompt = (
        "你是一个个人记忆系统的 AI 分析师。根据以下数据，写一段 3-5 句话的'本周洞察'（中文），"
        "像一位了解你的朋友一样，点出值得注意的模式、增长或盲点。\n\n"
        f"数据:\n- 总实体数: {entity_total}\n"
        f"- 实体类型分布: {type_counts}\n"
        f"- Top 实体: {top_entities}\n"
        f"- 记忆边数: {memory_count}\n"
        f"- 知识条目: {knowledge_count}\n"
        f"- 社交关系数: {social_stats.get('total_edges', 0)}\n"
        f"- 待重连: {reconnect_count}\n\n"
        "用友好、有洞察力的语气（不要列表，直接写段落文字）。"
        "如果数据很少（实体<10），就鼓励用户多聊。"
    )

    try:
        provider, model = _get_model_config()
        from agent.auxiliary_client import call_llm
        resp = call_llm(
            provider=provider, model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7, max_tokens=400, timeout=30,
        )
        insights_text = resp.choices[0].message.content.strip()
        return {"insights_text": insights_text, "generated_at": datetime.now().isoformat(), "data_snapshot": data_snapshot}
    except Exception as e:
        logger.warning("Insights generation failed: %s", e)
        return {"insights_text": "暂无洞察数据。多聊几句，我就能为你分析记忆模式。", "generated_at": datetime.now().isoformat(), "data_snapshot": data_snapshot}


@app.post("/api/report/monthly")
async def api_report_monthly():
    """Generate a monthly report with data aggregation + AI narrative."""
    hm = get_heimdall()
    gs = hm.provider.store
    sg = hm.provider.social_graph
    km = hm.provider.knowledge

    entity_total = gs.get_entity_count()
    memory_count = 0
    try:
        row = gs._conn.execute("SELECT COUNT(*) as cnt FROM heimdall_memory_edges").fetchone()
        if row:
            memory_count = row["cnt"]
    except Exception:
        pass

    knowledge_count = len(km.search("", limit=1000)) if km else 0
    social_stats = sg.get_stats() if sg else {}

    # Top entities this session
    top_entities = []
    try:
        cur = gs._conn.execute(
            "SELECT display_name, entity_type, occurrence_count FROM heimdall_entities "
            "WHERE status='active' ORDER BY occurrence_count DESC LIMIT 10"
        )
        top_entities = [dict(r) for r in cur.fetchall()]
    except Exception:
        pass

    # Emotion distribution
    emotion_dist = {"positive": 0, "neutral": 0, "negative": 0}
    try:
        cur = gs._conn.execute("SELECT emotion FROM heimdall_memory_edges WHERE emotion IS NOT NULL")
        for row in cur:
            v = row["emotion"]
            if v > 0.3:
                emotion_dist["positive"] += 1
            elif v < -0.3:
                emotion_dist["negative"] += 1
            else:
                emotion_dist["neutral"] += 1
    except Exception:
        pass

    # Persona drift data
    pm = hm.provider.persona
    drift_data = _persona_snapshot.get_drift_data(pm) if pm else {"snapshots": [], "current_drift": {}}

    # Pivot moments (Phase 2)
    pivot_moments = []
    try:
        from heimdall.views.pivots import PivotDetector
        detector = PivotDetector(gs)
        pivot_moments = detector.detect_for_summary()
    except Exception:
        pass

    today_str = datetime.now().strftime("%Y年%m月%d日")
    metrics = {
        "entity_total": entity_total,
        "memory_count": memory_count,
        "knowledge_count": knowledge_count,
        "social_edges": social_stats.get("total_edges", 0),
        "top_entities": top_entities,
        "emotion_distribution": emotion_dist,
        "drift": drift_data.get("current_drift", {}),
        "pivot_moments": pivot_moments,
    }

    # Build pivot moments summary for prompt
    pivots_text = ""
    if pivot_moments:
        pivots_text = "### 枢纽时刻\n"
        for i, p in enumerate(pivot_moments[:3]):
            pivots_text += f"- {p['date']}: {p['narrative']} (score={p['score']:.2f})\n"

    prompt = (
        f"生成一份个人记忆系统月度报告（{today_str}）。用中文 Markdown 格式，以下结构：\n\n"
        "## 月度记忆报告\n\n"
        "### 总览\n- 用1-2句话总结本月记忆积累情况\n\n"
        "### 关键数字\n- 实体总数、记忆条数、知识条目、社交关系数\n\n"
        "### 核心人物/项目\n- 列出 top 实体并简要描述其角色\n\n"
        "### 情绪概览\n- 正面/中性/负面记忆分布及趋势解读\n\n"
        "### 枢纽时刻\n- 基于检测到的枢纽节点，描述重要的转折或连接时刻\n\n"
        "### 人格稳定性\n- drift 状态评估\n\n"
        "### 下月建议\n- 3条具体的关注/改进建议\n\n"
        f"数据:\n- 实体总数: {entity_total}\n- 记忆条数: {memory_count}\n"
        f"- 知识条目: {knowledge_count}\n- 社交关系: {social_stats.get('total_edges', 0)}\n"
        f"- Top 实体: {top_entities}\n- 情绪分布: {emotion_dist}\n"
        f"- 人格漂移: {drift_data.get('current_drift', {})}\n"
        f"{pivots_text}\n"
    )

    try:
        provider, model = _get_model_config()
        from agent.auxiliary_client import call_llm
        resp = call_llm(
            provider=provider, model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7, max_tokens=1200, timeout=45,
        )
        report_text = resp.choices[0].message.content.strip()
        return {"report_text": report_text, "report_date": today_str, "metrics": metrics}
    except Exception as e:
        logger.warning("Monthly report generation failed: %s", e)
        return {"report_text": f"## 月度记忆报告\n\n报告生成失败: {e}\n\n请稍后重试。", "report_date": today_str, "metrics": metrics}


# ---------------------------------------------------------------------------
# Chat (SSE streaming)
# ---------------------------------------------------------------------------

@app.get("/api/chat/history")
async def api_chat_history():
    """Return the current session's conversation history."""
    return {"messages": _chat_history}


@app.post("/api/chat/stream")
async def api_chat_stream(request: Request):
    body = await request.json()
    user_message = body.get("message", "").strip()
    if not user_message:
        return StreamingResponse(_sse_stream("请输入消息"), media_type="text/event-stream")

    agent = get_agent()
    return StreamingResponse(
        _run_chat_stream(agent, user_message),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


async def _run_chat_stream(agent, user_message: str):
    """Run AIAgent.run_conversation with SSE output."""
    import queue
    import threading

    q: queue.Queue = queue.Queue()

    def _collect(events):
        collected = []
        for event in events:
            q.put(event)
            collected.append(event)
        return collected

    stream_callback = _collect

    # Track the response
    result = {"text": "", "done": False, "error": ""}

    # Snapshot previous history as context for the AI
    _history_snapshot = list(_chat_history)

    def _run():
        try:
            resp = agent.run_conversation(
                user_message=user_message,
                stream_callback=stream_callback,
                conversation_history=_history_snapshot if _history_snapshot else None,
            )
            result["text"] = resp
            result["done"] = True
        except Exception as e:
            result["error"] = str(e)
            result["done"] = True
        q.put(None)

    t = threading.Thread(target=_run, daemon=True)
    t.start()

    while True:
        try:
            item = q.get(timeout=0.1)
        except queue.Empty:
            # Check if done
            if result["done"] and q.empty():
                break
            await asyncio.sleep(0.05)
            continue

        if item is None:
            break

        if isinstance(item, dict) and item.get("type") == "delta":
            yield f"data: {json.dumps({'type': 'delta', 'text': item.get('text', '')})}\n\n"
        elif isinstance(item, str):
            yield f"data: {json.dumps({'type': 'delta', 'text': item})}\n\n"

    if result["error"]:
        yield f"data: {json.dumps({'type': 'error', 'text': result['error']})}\n\n"

    # Extract token info and reasoning from the response
    resp_data = result["text"] if isinstance(result["text"], dict) else {}
    done_payload = {"type": "done"}
    if resp_data:
        done_payload["input_tokens"] = resp_data.get("input_tokens", 0)
        done_payload["output_tokens"] = resp_data.get("output_tokens", 0)
        done_payload["total_tokens"] = resp_data.get("total_tokens", 0)
        done_payload["model"] = resp_data.get("model", "")
        done_payload["reasoning"] = resp_data.get("last_reasoning", "")
    yield f"data: {json.dumps(done_payload)}\n\n"

    # Save to history & sync to HEIMDALL (background — don't block stream close)
    assistant_text = resp_data.get("final_response", "") if resp_data else ""
    if not assistant_text and isinstance(result["text"], str):
        assistant_text = result["text"]
    _background_sync(user_message, assistant_text)


async def _sse_stream(text: str):
    yield f"data: {json.dumps({'type': 'delta', 'text': text})}\n\n"
    yield f"data: {json.dumps({'type': 'done'})}\n\n"


def _background_sync(user_msg: str, assistant_msg: str):
    """Save to chat history and run HEIMDALL sync in a background thread.

    Keeps the SSE stream from blocking on entity extraction / DB writes
    after the 'done' event has already been sent to the client.
    """
    import threading

    if not assistant_msg:
        return

    _chat_history.append({"role": "user", "content": user_msg})
    _chat_history.append({"role": "assistant", "content": assistant_msg})

    def _sync():
        try:
            hm = get_heimdall()
            hm.sync_turn(user_msg, assistant_msg)
        except Exception:
            pass
        _maybe_auto_snapshot()

    t = threading.Thread(target=_sync, daemon=True)
    t.start()


# ---------------------------------------------------------------------------
# Voice Chat — audio upload + transcription + streaming response
# ---------------------------------------------------------------------------

_VOICE_MODEL = None


def _get_voice_model():
    """Lazy-load faster-whisper model (base, CPU, int8)."""
    global _VOICE_MODEL
    if _VOICE_MODEL is None:
        import faster_whisper
        _VOICE_MODEL = faster_whisper.WhisperModel("base", device="cpu", compute_type="int8")
    return _VOICE_MODEL


@app.post("/api/chat/voice")
async def api_chat_voice(file: UploadFile = File(...)):
    """Receive a voice message (webm/wav/mp3), transcribe, and stream chat response."""
    if not file.filename:
        raise HTTPException(400, "缺少音频文件")
    ext = Path(file.filename).suffix.lower()
    if ext not in (".webm", ".wav", ".mp3", ".ogg", ".m4a", ".opus"):
        raise HTTPException(400, f"不支持的音频格式: {ext}")

    raw = await file.read()
    if len(raw) < 512:
        raise HTTPException(400, "音频文件太小，可能没有内容")

    import tempfile
    import subprocess

    # Save uploaded audio to temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp_in:
        tmp_in.write(raw)
        tmp_in_path = tmp_in.name

    wav_path = tmp_in_path + ".wav"

    try:
        # Convert to 16kHz mono WAV for whisper
        subprocess.run(
            ["ffmpeg", "-y", "-i", tmp_in_path, "-ar", "16000", "-ac", "1",
             "-sample_fmt", "s16", wav_path],
            capture_output=True, timeout=30,
            check=True,
        )

        model = _get_voice_model()
        segments, _info = model.transcribe(wav_path, language="zh")
        transcription = " ".join(seg.text.strip() for seg in segments)

        if not transcription:
            transcription = "（未识别到语音内容）"
    except subprocess.CalledProcessError as e:
        logger.warning("ffmpeg conversion failed: %s", e.stderr.decode() if e.stderr else str(e))
        raise HTTPException(400, "音频转换失败，请确认文件格式正确")
    except Exception as e:
        logger.warning("Voice transcription failed: %s", e)
        raise HTTPException(500, f"语音识别失败: {e}")
    finally:
        # Cleanup temp files
        for p in (tmp_in_path, wav_path):
            try:
                os.unlink(p)
            except OSError:
                pass

    agent = get_agent()
    return StreamingResponse(
        _run_voice_chat_stream(agent, transcription, raw, ext),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


async def _run_voice_chat_stream(agent, transcription: str, audio_raw: bytes, ext: str):
    """Stream chat response, starting with the transcription and audio blob."""
    import base64
    import queue
    import threading

    # 1. Emit transcription + audio data for the frontend voice bubble
    audio_b64 = base64.b64encode(audio_raw).decode("ascii")
    voice_event = {
        "type": "voice_transcription",
        "text": transcription,
        "audio_base64": audio_b64,
        "audio_mime": f"audio/{ext.lstrip('.')}",
    }
    yield f"data: {json.dumps(voice_event)}\n\n"

    # 2. Run normal chat pipeline
    q: queue.Queue = queue.Queue()

    def _collect(events):
        for event in events:
            q.put(event)
        return []

    result = {"text": "", "done": False, "error": ""}
    _history_snapshot = list(_chat_history)

    def _run():
        try:
            resp = agent.run_conversation(
                user_message=transcription,
                stream_callback=_collect,
                conversation_history=_history_snapshot if _history_snapshot else None,
            )
            result["text"] = resp
            result["done"] = True
        except Exception as e:
            result["error"] = str(e)
            result["done"] = True
        q.put(None)

    t = threading.Thread(target=_run, daemon=True)
    t.start()

    while True:
        try:
            item = q.get(timeout=0.1)
        except queue.Empty:
            if result["done"] and q.empty():
                break
            await asyncio.sleep(0.05)
            continue
        if item is None:
            break
        if isinstance(item, dict) and item.get("type") == "delta":
            yield f"data: {json.dumps({'type': 'delta', 'text': item.get('text', '')})}\n\n"
        elif isinstance(item, str):
            yield f"data: {json.dumps({'type': 'delta', 'text': item})}\n\n"

    if result["error"]:
        yield f"data: {json.dumps({'type': 'error', 'text': result['error']})}\n\n"

    resp_data = result["text"] if isinstance(result["text"], dict) else {}
    done_payload = {"type": "done"}
    if resp_data:
        done_payload["input_tokens"] = resp_data.get("input_tokens", 0)
        done_payload["output_tokens"] = resp_data.get("output_tokens", 0)
        done_payload["total_tokens"] = resp_data.get("total_tokens", 0)
        done_payload["model"] = resp_data.get("model", "")
        done_payload["reasoning"] = resp_data.get("last_reasoning", "")
    yield f"data: {json.dumps(done_payload)}\n\n"

    assistant_text = resp_data.get("final_response", "") if resp_data else ""
    if not assistant_text and isinstance(result["text"], str):
        assistant_text = result["text"]

    _background_sync(transcription, assistant_text)


# ---------------------------------------------------------------------------
# File Browser
# ---------------------------------------------------------------------------

_FILE_WHITELIST = {
    "persona.md": ("PERSONA.md", "📛 人格定义"),
    "knowledge.md": ("KNOWLEDGE.md", "📚 知识域"),
    "config.yaml": ("config.yaml", "⚙️ 配置"),
    "soul.md": ("SOUL.md", "🧠 系统灵魂"),
    "agent.log": ("logs/agent.log", "📋 事件日志"),
    "snapshots.json": ("persona_snapshots.json", "📉 人格快照"),
}


@app.get("/api/files")
async def api_files():
    """List browsable files."""
    hermes_home = _get_hermes_home()
    files = []
    for key, (rel_path, desc) in _FILE_WHITELIST.items():
        fpath = hermes_home / rel_path
        exists = fpath.exists()
        size = fpath.stat().st_size if exists else 0
        files.append({
            "key": key,
            "name": rel_path.split("/")[-1],
            "description": desc,
            "path": str(fpath),
            "exists": exists,
            "size": size,
        })
    return {"files": files}


@app.get("/api/files/{name}")
async def api_file_read(name: str):
    """Read a whitelisted file's content."""
    if name not in _FILE_WHITELIST:
        raise HTTPException(403, "File not in whitelist")
    rel_path = _FILE_WHITELIST[name][0]
    fpath = _get_hermes_home() / rel_path
    if not fpath.exists():
        raise HTTPException(404, f"File not found: {rel_path}")
    content = fpath.read_text("utf-8")
    return {"name": rel_path.split("/")[-1], "path": str(fpath), "content": content, "size": len(content)}


# ---------------------------------------------------------------------------
# File upload — preview + entity extraction
# ---------------------------------------------------------------------------

ALLOWED_EXTENSIONS = {".txt", ".md", ".json", ".yaml", ".yml", ".py", ".log", ".csv"}
MAX_UPLOAD_BYTES = 2 * 1024 * 1024  # 2 MB


@app.post("/api/upload/preview")
async def api_upload_preview(file: UploadFile = File(...)):
    """Preview an uploaded text file and extract entities from its content."""
    ext = Path(file.filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"不支持的文件类型: {ext}。支持: {', '.join(sorted(ALLOWED_EXTENSIONS))}")
    size = file.size or 0
    if size > MAX_UPLOAD_BYTES:
        raise HTTPException(400, f"文件过大 ({size} bytes)，限制 {MAX_UPLOAD_BYTES} bytes")
    try:
        raw = await file.read()
        content = raw.decode("utf-8")
    except UnicodeDecodeError:
        raise HTTPException(400, "无法以 UTF-8 解码文件内容")
    preview = content[:5000]
    entities_preview = []
    try:
        hm = get_heimdall()
        extractor = hm.extractor if hasattr(hm, "extractor") else None
        if extractor and hasattr(extractor, "extract_from_text"):
            result = extractor.extract_from_text(preview, session_id="upload-preview")
            entities_preview = result.get("entities", [])[:10]
    except Exception:
        pass
    return {
        "filename": file.filename,
        "size": size,
        "preview": preview,
        "total_chars": len(content),
        "entities_preview": entities_preview,
    }


@app.get("/api/db/tables")
async def api_db_tables():
    """Return SQLite table stats."""
    hm = get_heimdall()
    gs = hm.provider.store
    tables = []
    if gs._conn:
        try:
            rows = gs._conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' AND name LIKE 'heimdall_%' ORDER BY name"
            ).fetchall()
            for r in rows:
                tname = r["name"]
                cnt = gs._conn.execute(f"SELECT COUNT(*) as c FROM [{tname}]").fetchone()
                tables.append({"table": tname, "rows": cnt["c"] if cnt else 0})
        except Exception:
            pass
    # DB file size
    db_path = hm.provider.heimdall_dir / "heimdall.db" if hasattr(hm.provider, "heimdall_dir") else None
    db_size = db_path.stat().st_size if db_path and db_path.exists() else 0
    return {"tables": tables, "db_size": db_size, "db_path": str(db_path) if db_path else ""}


# ---------------------------------------------------------------------------
# Tools
# ---------------------------------------------------------------------------

@app.get("/api/tools")
async def api_tools():
    hm = get_heimdall()
    return {"tools": hm.get_tool_schemas()}


# ---------------------------------------------------------------------------
# Knowledge Ring V1.0 API endpoints
# ---------------------------------------------------------------------------

@app.get("/api/v1/daily-report")
async def api_v1_daily_report(date: str = ""):
    """Daily growth report from Knowledge Ring event_log."""
    get_heimdall()  # ensure engine is initialized
    from heimdall.views.daily import get_daily_report as _get_daily
    report = _get_daily(date or None)
    return report


@app.get("/api/v1/ring")
async def api_v1_ring():
    """Full Knowledge Ring graph data (entities + relations) for D3.js visualization."""
    hm = get_heimdall()
    data = hm.provider.store.get_ring_graph_data()
    return data


@app.get("/api/v1/entity/{entity_id}")
async def api_v1_entity_detail(entity_id: str):
    """Entity detail card: entity + aliases + relations + timeline."""
    hm = get_heimdall()
    entity = hm.provider.store.get_entity_v2(entity_id)
    if not entity:
        raise HTTPException(status_code=404, detail="Entity not found")
    aliases = hm.provider.store.get_aliases(entity_id)
    relations = hm.provider.store.get_relations(entity_id)
    return {
        "entity": entity,
        "aliases": aliases,
        "relations": relations,
    }


@app.get("/api/v1/domains")
async def api_v1_domains():
    """List all registered knowledge domains with stats."""
    hm = get_heimdall()
    domains = hm.provider.store.get_domains()
    stats = hm.provider.store.get_domain_stats()
    return {"domains": domains, "stats": stats}


@app.post("/api/v1/entity/{entity_id}/correct")
async def api_v1_entity_correct(entity_id: str, request: Request):
    """User correction: edit entity fields (name, type_detail, domains, properties)."""
    hm = get_heimdall()
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")

    allowed = {"name", "type_detail", "domains", "properties", "confidence"}
    updates = {k: v for k, v in body.items() if k in allowed}
    if not updates:
        raise HTTPException(status_code=400, detail="No valid fields to update")

    success = hm.provider.store.update_entity_v2(entity_id, **updates)
    if not success:
        raise HTTPException(status_code=404, detail="Entity not found")

    # Log as user correction in event_log
    entity = hm.provider.store.get_entity_v2(entity_id)
    return {"status": "corrected", "entity": entity}


@app.get("/api/v1/timeline")
async def api_v1_entity_timeline(entity_id: str = ""):
    """Get timeline events for an entity from event_log."""
    hm = get_heimdall()
    store = hm.provider.store
    if not store._conn:
        return {"events": []}
    if entity_id:
        rows = store._conn.execute(
            "SELECT * FROM kr_event_log WHERE entity_id = ? ORDER BY timestamp DESC LIMIT 50",
            (entity_id,),
        ).fetchall()
    else:
        rows = store._conn.execute(
            "SELECT * FROM kr_event_log ORDER BY timestamp DESC LIMIT 50"
        ).fetchall()
    return {"events": [dict(r) for r in rows]}


# ---------------------------------------------------------------------------
# V2.3 Knowledge API endpoints
# ---------------------------------------------------------------------------

@app.get("/api/knowledge/namespaces")
async def api_knowledge_namespaces():
    """List all knowledge namespaces with entity counts."""
    hm = get_heimdall()
    store = hm.provider.store
    try:
        namespaces = store.ops.list_namespaces() if store.ops else []
    except Exception:
        namespaces = []
    return {"namespaces": namespaces}


@app.get("/api/knowledge/graph-data")
async def api_knowledge_graph_data(namespace: str = None):
    """Return {nodes, edges} from kr_entities/kr_relations with optional namespace filter."""
    hm = get_heimdall()
    store = hm.provider.store
    data = store.get_ring_graph_data()

    # Map engine field names to desktop field names
    nodes = []
    for n in data.get("nodes", []):
        types_list = n.get("types", ["concept"])
        nodes.append({
            "id": n.get("entity_id", ""),
            "name": n.get("name", ""),
            "type": types_list[0] if types_list else "concept",
            "types": types_list,
            "description": n.get("description", n.get("type_detail", "")),
            "importance": n.get("confidence", 0.5),
            "degree": 0,
            "namespace": n.get("namespace", "general"),
        })

    edges = []
    for e in data.get("edges", []):
        edges.append({
            "source": e.get("source_id", ""),
            "target": e.get("target_id", ""),
            "type": e.get("type", ""),
            "weight": e.get("confidence", 0.7),
            "bidirectional": e.get("direction", "") == "bidirectional",
        })

    if namespace and namespace != "all":
        nodes = [n for n in nodes if n.get("namespace", "general") == namespace]
        node_ids = {n["id"] for n in nodes}
        edges = [
            e for e in edges
            if e["source"] in node_ids and e["target"] in node_ids
        ]

    return {"nodes": nodes, "edges": edges}


@app.get("/api/knowledge/neighbors")
async def api_knowledge_neighbors(entity_id: Optional[str] = None, hops: int = 2, namespace: Optional[str] = None, id: Optional[str] = None):
    """BFS neighbor query from kr_entities/kr_relations."""
    eid = entity_id or id  # accept both entity_id (engine) and id (desktop client)
    if not eid:
        return {"nodes": [], "edges": []}
    hm = get_heimdall()
    store = hm.provider.store
    if not store._conn:
        return {"nodes": [], "edges": []}

    visited: set[str] = {eid}
    frontier = {eid}
    all_edges: list[dict] = []

    for _ in range(max(1, min(hops, 5))):
        next_frontier: set[str] = set()
        for eid in frontier:
            relations = store.get_relations(eid)
            for r in relations:
                neighbor = r["target_id"] if r["source_id"] == eid else r["source_id"]
                if neighbor not in visited:
                    visited.add(neighbor)
                    next_frontier.add(neighbor)
                all_edges.append(r)
        frontier = next_frontier
        if not frontier:
            break

    nodes = []
    for eid in visited:
        entity = store.get_entity_v2(eid)
        if entity:
            if namespace and entity.get("namespace", "general") != namespace:
                continue
            types_list = entity.get("types", entity.get("type", ["concept"]))
            if isinstance(types_list, str):
                try:
                    import json as _json
                    types_list = _json.loads(types_list)
                except Exception:
                    types_list = ["concept"]
            nodes.append({
                "id": entity.get("entity_id", entity.get("id", "")),
                "name": entity.get("name", ""),
                "type": types_list[0] if types_list else "concept",
                "description": entity.get("description", entity.get("type_detail", "")),
                "importance": entity.get("confidence", entity.get("importance", 0.5)),
                "degree": 0,
                "namespace": entity.get("namespace", "general"),
            })

    mapped_edges = []
    edge_seen = set()
    for e in all_edges:
        key = (e.get("source_id", ""), e.get("target_id", ""))
        if key in edge_seen:
            continue
        edge_seen.add(key)
        mapped_edges.append({
            "source": e.get("source_id", ""),
            "target": e.get("target_id", ""),
            "type": e.get("type", ""),
            "weight": e.get("confidence", 0.7),
            "bidirectional": e.get("direction", "") == "bidirectional",
        })

    return {"nodes": nodes, "edges": mapped_edges}


@app.get("/api/knowledge/paths")
async def api_knowledge_paths(source_id: Optional[str] = None, target_id: Optional[str] = None, max_len: int = 5,
                               source: Optional[str] = None, target: Optional[str] = None):
    """BFS path finding between two entities."""
    sid = source_id or source  # accept both engine and desktop client param names
    tid = target_id or target
    if not sid or not tid:
        return {"paths": [], "count": 0}
    hm = get_heimdall()
    store = hm.provider.store
    paths = store.find_relation_path(sid, tid, max_len)
    return {"paths": paths, "count": len(paths)}


@app.get("/api/knowledge/stats")
async def api_knowledge_stats(namespace: str = None):
    """Return {totalNodes, totalEdges, avgDegree, typeDistribution}."""
    hm = get_heimdall()
    store = hm.provider.store
    total_nodes = store.ops.count_by_namespace(namespace) if (namespace and store.ops) else (store.get_entity_count() if store else 0)
    total_edges = store.get_relation_count() if store else 0
    avg_degree = round((2 * total_edges / total_nodes), 2) if total_nodes > 0 else 0
    domain_stats = store.ops.get_domain_stats() if store.ops else []
    return {
        "totalNodes": total_nodes,
        "totalEdges": total_edges,
        "avgDegree": avg_degree,
        "typeDistribution": domain_stats,
    }


@app.get("/api/knowledge/operation-log")
async def api_knowledge_operation_log(limit: int = 50, namespace: str = None):
    """Recent kr_event_log entries."""
    hm = get_heimdall()
    store = hm.provider.store
    if not store._conn:
        return {"events": []}
    rows = store._conn.execute(
        "SELECT * FROM kr_event_log ORDER BY timestamp DESC LIMIT ?", (limit,)
    ).fetchall()
    return {"events": [dict(r) for r in rows]}


@app.get("/api/knowledge/daily-digest")
async def api_knowledge_daily_digest(days: int = 7, namespace: str = None):
    """Daily operation summary for the last N days."""
    from datetime import date, timedelta
    hm = get_heimdall()
    store = hm.provider.store
    if not store._conn:
        return {"digest": []}

    today = date.today()
    digest = []
    for i in range(days):
        d = (today - timedelta(days=i)).isoformat()
        events = store.ops.get_daily_events(d) if store.ops else []
        digest.append({
            "date": d,
            "event_count": len(events),
            "top_types": list(set(e.get("event_type", "") for e in events[:5])),
        })
    return {"digest": digest}


@app.get("/api/knowledge/events/stream")
async def api_knowledge_events_stream(namespace: str = "general", client_id: str = ""):
    """SSE stream of knowledge graph change events for desktop live updates."""
    from sqlite3 import OperationalError as _SqliteOpError

    async def event_generator():
        hm = get_heimdall()
        store = hm.provider.store
        last_id = 0
        # Seed from the most recent event so we only emit new ones
        if store._conn:
            try:
                row = store._conn.execute(
                    "SELECT MAX(log_id) as max_id FROM kr_event_log"
                ).fetchone()
                if row and row["max_id"]:
                    last_id = int(row["max_id"])
            except (_SqliteOpError, Exception):
                pass

        heartbeat_interval = 30
        last_heartbeat = time.time()

        while True:
            if await asyncio.sleep(2, result=True) is False:
                break

            now = time.time()
            if now - last_heartbeat >= heartbeat_interval:
                last_heartbeat = now
                yield ": heartbeat\n\n"
                continue

            if not store._conn:
                continue

            try:
                rows = store._conn.execute(
                    "SELECT * FROM kr_event_log WHERE log_id > ? ORDER BY log_id ASC LIMIT 50",
                    (last_id,),
                ).fetchall()
                for row in rows:
                    row_dict = dict(row)
                    event_type = row_dict.get("event_type", "unknown")
                    payload = json.dumps(row_dict, default=str)
                    yield f"event: {event_type}\ndata: {payload}\n\n"
                    if row_dict["log_id"] > last_id:
                        last_id = row_dict["log_id"]
            except (_SqliteOpError, Exception):
                continue

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/api/extract")
async def api_extract(request: Request):
    """Entity extraction for desktop auto-extraction flow.

    Supports two modes:
      - mode=entities: extract entities from raw text (trigger-word extraction)
      - mode=auto:    extract from user+assistant conversation turn
    """
    hm = get_heimdall()
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")

    mode = body.get("mode", "auto")

    if mode == "entities":
        text = body.get("text", "")
        if not text or len(text) < 20:
            return {"entities": []}
        try:
            result = hm.provider.extractor.extract_from_text(
                text, session_id="desktop-extract", namespace=body.get("namespace", "general")
            )
        except Exception:
            return {"entities": []}
        return {"entities": result.get("entities", [])}

    # mode == "auto": conversation turn extraction
    user_msg = body.get("user_message", "")
    asst_msg = body.get("assistant_response", "")
    if not user_msg and not asst_msg:
        return {"extracted": 0, "reasoning": "Empty input"}

    try:
        result = hm.provider.extractor.extract_from_turn(
            user_message=user_msg,
            assistant_message=asst_msg,
            session_id=body.get("session_id", "desktop-auto"),
            namespace=body.get("namespace", "general"),
        )
    except Exception as e:
        return {"extracted": 0, "reasoning": f"Extraction failed: {str(e)}"}

    return {
        "extracted": result.get("extracted_count", 0),
        "reasoning": result.get("reasoning", ""),
    }


@app.post("/api/knowledge/sync-wiki")
async def api_knowledge_sync_wiki(request: Request):
    """Sync a wiki document: queue for async extraction, return immediately.

    The desktop already does local extraction via importWikiFileLocally()
    for instant visibility. Engine extraction is an enhancement that runs
    in the background via pending_sync queue + background worker.
    """
    hm = get_heimdall()
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")

    doc_path = body.get("path", "")
    content = body.get("content", "")
    namespace = body.get("namespace", "general")

    if not content.strip():
        return {"synced": False, "error": "empty content"}

    # Queue for async extraction — never block the request thread on LLM calls
    if hm.provider.store._conn:
        try:
            import uuid as _uuid
            hm.provider.store._conn.execute(
                "INSERT INTO kr_pending_sync (id, payload_type, payload_json, namespace) "
                "VALUES (?, 'wiki', ?, ?)",
                (_uuid.uuid4().hex, json.dumps(body, ensure_ascii=False), namespace),
            )
            return {"synced": False, "queued": True, "message": "queued for background extraction"}
        except Exception as e:
            return {"synced": False, "queued": True, "error": str(e)}

    return {"synced": False, "queued": False, "error": "store unavailable"}


@app.post("/api/knowledge/sync-canvas")
async def api_knowledge_sync_canvas(request: Request):
    """Sync a canvas document: write nodes/edges to heimdall.db."""
    hm = get_heimdall()
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")

    namespace = body.get("namespace", "general")
    nodes = body.get("data", {}).get("nodes", [])
    edges = body.get("data", {}).get("edges", [])

    entity_count = 0
    relation_count = 0

    try:
        for node in nodes:
            node_type = node.get("_hermes", {}).get("entityType", "concept")
            hm.provider.store.upsert_entity_v2(
                name=node.get("label", node.get("id", "")),
                entity_type=node_type,
                properties={"canvas_id": node.get("id", ""), "description": node.get("description", "")},
                namespace=namespace,
            )
            entity_count += 1

        for edge in edges:
            src = edge.get("source", "")
            tgt = edge.get("target", "")
            if src and tgt:
                src_entity = hm.provider.store.get_entity_v2_by_name(src)
                tgt_entity = hm.provider.store.get_entity_v2_by_name(tgt)
                if not src_entity:
                    eid = hm.provider.store.upsert_entity_v2(name=src, namespace=namespace)
                    src_entity = hm.provider.store.get_entity_v2(eid)
                if not tgt_entity:
                    eid = hm.provider.store.upsert_entity_v2(name=tgt, namespace=namespace)
                    tgt_entity = hm.provider.store.get_entity_v2(eid)
                if src_entity and tgt_entity:
                    hm.provider.store.add_relation(
                        source_id=src_entity["entity_id"],
                        target_id=tgt_entity["entity_id"],
                        rel_type=edge.get("label", "relates_to"),
                        namespace=namespace,
                    )
                    relation_count += 1
    except Exception as e:
        if hm.provider.store._conn:
            import uuid as _uuid
            hm.provider.store._conn.execute(
                "INSERT INTO kr_pending_sync (id, payload_type, payload_json, namespace) "
                "VALUES (?, 'canvas', ?, ?)",
                (_uuid.uuid4().hex, json.dumps(body, ensure_ascii=False), namespace),
            )
        return {"synced": False, "queued": True, "error": str(e)}

    # Trigger overview rebuild
    if hm.provider._overview_cache and namespace:
        try:
            hm.provider._overview_cache.schedule_rebuild(namespace)
        except Exception:
            pass

    return {"synced": True, "entity_count": entity_count, "relation_count": relation_count}


@app.get("/api/knowledge/pending-failures")
async def api_list_pending_failures(namespace: str = None):
    """List failed pending sync records (retry_count >= 5) for UI visibility."""
    hm = get_heimdall()
    failures = hm.provider.store.get_pending_failures(namespace)
    return {"failures": failures}


@app.post("/api/knowledge/pending-cleanup")
async def api_cleanup_pending(days: int = 7):
    """Manually purge dead pending sync rows older than N days."""
    hm = get_heimdall()
    deleted = hm.provider.store.cleanup_dead_pending_sync(days)
    return {"deleted": deleted}


@app.delete("/api/knowledge/pending/{pending_id}")
async def api_delete_pending(pending_id: str):
    """Discard a specific pending sync record."""
    hm = get_heimdall()
    hm.provider.store.delete_pending_sync(pending_id)
    return {"deleted": True, "id": pending_id}


@app.get("/api/knowledge/pending-stats")
async def api_pending_stats():
    """Return counts of pending + failed sync records."""
    hm = get_heimdall()
    conn = hm.provider.store._conn
    pending = conn.execute(
        "SELECT COUNT(*) FROM kr_pending_sync WHERE retry_count < 5"
    ).fetchone()[0]
    failed = conn.execute(
        "SELECT COUNT(*) FROM kr_pending_sync WHERE retry_count >= 5"
    ).fetchone()[0]
    return {"pending": pending, "failed": failed}


@app.post("/api/convert")
async def api_convert_document(request: Request):
    """Convert .docx/.xlsx/.pptx to markdown (pure Python, no LLM)."""
    import base64
    import io

    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")

    ext = body.get("extension", "").lower().lstrip(".")
    b64 = body.get("content", "")

    if not b64 or ext not in ("docx", "xlsx", "pptx"):
        raise HTTPException(status_code=400, detail="extension must be docx, xlsx, or pptx; content required")

    try:
        raw = base64.b64decode(b64)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid base64 content")

    md = ""

    if ext == "docx":
        try:
            from docx import Document
        except ImportError:
            raise HTTPException(status_code=500, detail="python-docx not installed")
        doc = Document(io.BytesIO(raw))
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                parts.append("")
                continue
            if para.style.name.startswith("Heading"):
                level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
        # Tables
        for table in doc.tables:
            parts.append("")
            for ri, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                parts.append("| " + " | ".join(cells) + " |")
                if ri == 0:
                    parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
            parts.append("")
        md = "\n\n".join(parts)

    elif ext == "xlsx":
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise HTTPException(status_code=500, detail="openpyxl not installed")
        wb = load_workbook(io.BytesIO(raw), data_only=True)
        parts: list[str] = []
        for sname in wb.sheetnames:
            ws = wb[sname]
            parts.append(f"## {sname}")
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                parts.append("_(empty sheet)_")
                continue
            for ri, row in enumerate(rows):
                cells = [str(c or "").replace("\n", " ").replace("|", "\\|") for c in row]
                parts.append("| " + " | ".join(cells) + " |")
                if ri == 0:
                    parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
            parts.append("")
        md = "\n\n".join(parts)

    elif ext == "pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise HTTPException(status_code=500, detail="python-pptx not installed")
        prs = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        for si, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {si}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
            parts.append("")
        md = "\n\n".join(parts)

    return {"markdown": md}


# ---------------------------------------------------------------------------
# Phase 1-2: Inference + Causal Chain + Evolution endpoints
# ---------------------------------------------------------------------------


@app.post("/api/inference/run")
async def api_inference_run(request: Request):
    """Run transitive inference over a namespace to find implied relationships."""
    body = await request.json()
    namespace = body.get("namespace", "general")
    min_shared = body.get("min_shared_neighbors", 2)

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.inference import TransitiveInferenceEngine
    engine = TransitiveInferenceEngine(store)

    inferences = engine.run(namespace=namespace, min_shared_neighbors=min_shared)

    # Persist to database
    now = time.strftime("%Y-%m-%dT%H:%M:%S")
    for inf in inferences:
        try:
            store._conn.execute(
                """INSERT OR IGNORE INTO kr_inferences
                   (inference_id, entity_a, entity_b, inferred_type, evidence,
                    confidence, status, namespace, created_at)
                   VALUES (?, ?, ?, ?, ?, ?, 'pending', ?, ?)""",
                (inf.inference_id, inf.entity_a, inf.entity_b, inf.inferred_type,
                 json.dumps(inf.evidence, ensure_ascii=False),
                 inf.confidence, inf.namespace, now),
            )
        except Exception:
            pass
    store._conn.commit()

    return {
        "inferences": [
            {
                "inference_id": i.inference_id,
                "entity_a": i.entity_a,
                "entity_b": i.entity_b,
                "inferred_type": i.inferred_type,
                "evidence": i.evidence,
                "confidence": i.confidence,
                "status": i.status,
            }
            for i in inferences
        ],
        "count": len(inferences),
    }


@app.get("/api/inference/candidates")
async def api_inference_candidates(namespace: str = "general", limit: int = 10, status: str = "pending"):
    """List pending inference candidates."""
    hm = get_heimdall()
    store = hm.provider.store
    rows = store._conn.execute(
        """SELECT * FROM kr_inferences
           WHERE namespace = ? AND status = ?
           ORDER BY confidence DESC LIMIT ?""",
        (namespace, status, limit),
    ).fetchall()
    return {"inferences": [dict(r) for r in rows]}


@app.post("/api/inference/{inference_id}/confirm")
async def api_inference_confirm(inference_id: str):
    """Confirm a pending inference, promoting it to a real relation."""
    hm = get_heimdall()
    store = hm.provider.store

    row = store._conn.execute(
        "SELECT * FROM kr_inferences WHERE inference_id = ?", (inference_id,)
    ).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Inference not found")

    now = time.strftime("%Y-%m-%dT%H:%M:%S")
    store._conn.execute(
        "UPDATE kr_inferences SET status = 'confirmed', resolved_at = ? WHERE inference_id = ?",
        (now, inference_id),
    )
    # Create the actual relation
    store.ops.add_relation(
        source_id=row["entity_a"],
        target_id=row["entity_b"],
        rel_type=row["inferred_type"],
        confidence=0.7,
        namespace=row["namespace"],
    )
    store._conn.commit()
    return {"status": "confirmed", "inference_id": inference_id}


@app.post("/api/inference/{inference_id}/reject")
async def api_inference_reject(inference_id: str):
    """Reject a pending inference."""
    hm = get_heimdall()
    store = hm.provider.store

    now = time.strftime("%Y-%m-%dT%H:%M:%S")
    store._conn.execute(
        "UPDATE kr_inferences SET status = 'rejected', resolved_at = ? WHERE inference_id = ?",
        (now, inference_id),
    )
    store._conn.commit()
    return {"status": "rejected", "inference_id": inference_id}


# Causal chain endpoints

@app.post("/api/evolution/causal/rebuild")
async def api_causal_rebuild(request: Request):
    """Rebuild causal chains for a namespace from all causal relations."""
    body = await request.json()
    namespace = body.get("namespace", "general")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.causal import CausalChainBuilder
    builder = CausalChainBuilder(store)
    chains = builder.rebuild_all(namespace=namespace)

    return {
        "chains": [
            {
                "chain_id": c.chain_id,
                "length": c.length,
                "chain_score": round(c.chain_score, 4),
                "is_active": c.is_active,
            }
            for c in chains
        ],
        "count": len(chains),
    }


@app.get("/api/evolution/causal/chains")
async def api_causal_chains(namespace: str = "general", min_score: float = 0.3):
    """List significant causal chains for a namespace."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.causal import CausalChainBuilder
    builder = CausalChainBuilder(store)
    chains = builder.find_significant_chains(namespace=namespace, min_score=min_score)
    return {"chains": chains, "count": len(chains)}


# ---------------------------------------------------------------------------
# Phase 2: Conflict detection + Confidence evolution
# ---------------------------------------------------------------------------


@app.post("/api/evolution/conflicts/check")
async def api_conflict_check(request: Request):
    """Check for conflicts when new properties are submitted for an entity."""
    body = await request.json()
    entity_id = body.get("entity_id", "")
    properties = body.get("properties", {})
    namespace = body.get("namespace", "general")
    source = body.get("source", "manual")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.conflict import ConflictDetector
    detector = ConflictDetector(store)
    conflicts = detector.check(
        entity_id=entity_id,
        new_properties=properties,
        source=source,
        namespace=namespace,
    )

    # Persist conflicts
    for c in conflicts:
        detector.record_change(
            entity_id=c.entity_id,
            field=c.field,
            old_value=c.old_value,
            new_value=c.new_value,
            change_type=c.change_type,
            source=c.source,
            namespace=c.namespace,
        )

    return {
        "conflict_found": len(conflicts) > 0,
        "conflicts": [
            {
                "conflict_id": c.conflict_id,
                "entity_id": c.entity_id,
                "field": c.field,
                "old_value": c.old_value,
                "new_value": c.new_value,
                "change_type": c.change_type,
            }
            for c in conflicts
        ],
        "count": len(conflicts),
    }


@app.get("/api/evolution/conflicts")
async def api_conflict_list(namespace: str = "general", status: str = "pending"):
    """List conflict records for a namespace."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.conflict import ConflictDetector
    detector = ConflictDetector(store)
    conflicts = detector.scan_namespace(namespace=namespace)
    return {"conflicts": conflicts, "count": len(conflicts)}


@app.post("/api/evolution/conflicts/{history_id}/resolve")
async def api_conflict_resolve(history_id: str, request: Request):
    """Resolve a specific conflict by history_id."""
    body = await request.json()
    resolution = body.get("resolution", "dismiss")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.conflict import ConflictDetector
    detector = ConflictDetector(store)

    try:
        result = detector.resolve(history_id, resolution)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/evolution/confidence/boost")
async def api_confidence_boost(request: Request):
    """Manually boost an entity's confidence."""
    body = await request.json()
    entity_id = body.get("entity_id", "")
    amount = body.get("amount", 0.1)
    reason = body.get("reason", "manual_confirm")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.confidence import ConfidenceManager
    mgr = ConfidenceManager(store)
    result = mgr.boost(entity_id=entity_id, amount=amount, reason=reason)

    return {
        "entity_id": result.entity_id,
        "old_confidence": result.old_confidence,
        "new_confidence": result.new_confidence,
        "action": result.action,
    }


@app.post("/api/evolution/confidence/decay")
async def api_confidence_decay(request: Request):
    """Manually decay an entity's confidence."""
    body = await request.json()
    entity_id = body.get("entity_id", "")
    amount = body.get("amount", 0.05)
    reason = body.get("reason", "marked_outdated")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.confidence import ConfidenceManager
    mgr = ConfidenceManager(store)
    result = mgr.decay(entity_id=entity_id, amount=amount, reason=reason)

    return {
        "entity_id": result.entity_id,
        "old_confidence": result.old_confidence,
        "new_confidence": result.new_confidence,
        "action": result.action,
    }


@app.post("/api/evolution/confidence/apply-time-decay")
async def api_confidence_apply_time_decay(request: Request):
    """Apply exponential time decay to all entities in a namespace."""
    body = await request.json()
    namespace = body.get("namespace", "general")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.confidence import ConfidenceManager
    mgr = ConfidenceManager(store)
    results = mgr.apply_time_decay(namespace=namespace)

    archived = sum(1 for r in results if r.action == "archived")
    return {
        "entities_updated": len(results),
        "archived": archived,
    }


# ---------------------------------------------------------------------------
# Phase 3: Hypothesis generation + Viewpoint tracking
# ---------------------------------------------------------------------------


@app.get("/api/hypothesis/gaps")
async def api_hypothesis_gaps(namespace: str = "general", entity_type: str = None):
    """Find entities missing common relation types."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.hypothesis import HypothesisEngine
    engine = HypothesisEngine(store)
    gaps = engine.find_gaps(namespace=namespace, entity_type=entity_type)

    return {
        "gaps": [
            {
                "entity_id": g.entity_id,
                "entity_name": g.entity_name,
                "missing_relation_type": g.missing_relation_type,
                "suggestion": g.suggestion,
            }
            for g in gaps
        ],
        "count": len(gaps),
    }


@app.get("/api/hypothesis/generate")
async def api_hypothesis_generate(namespace: str = "general", limit: int = 20):
    """Generate hypotheses for entity pairs lacking relations."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.hypothesis import HypothesisEngine
    engine = HypothesisEngine(store)
    hypotheses = engine.generate(namespace=namespace, limit=limit)

    return {
        "hypotheses": [
            {
                "hypothesis_id": h.hypothesis_id,
                "entity_a": h.entity_a,
                "entity_b": h.entity_b,
                "suggested_relation": h.suggested_relation,
                "reasoning": h.reasoning,
                "confidence": h.confidence,
                "status": h.status,
            }
            for h in hypotheses
        ],
        "count": len(hypotheses),
    }


@app.post("/api/hypothesis/{hypothesis_id}/dismiss")
async def api_hypothesis_dismiss(hypothesis_id: str):
    """Dismiss a generated hypothesis."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.hypothesis import HypothesisEngine
    engine = HypothesisEngine(store)
    engine.dismiss(hypothesis_id)
    return {"status": "dismissed", "hypothesis_id": hypothesis_id}


@app.get("/api/evolution/viewpoint/drifted")
async def api_viewpoint_drifted(namespace: str = "general", min_changes: int = 2):
    """Find entities with significant viewpoint drift.

    IMPORTANT: This route MUST be registered before /api/evolution/viewpoint/{entity_id}
    to avoid FastAPI matching "drifted" as an {entity_id} parameter.
    """
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.viewpoint import ViewpointTracker
    tracker = ViewpointTracker(store)
    entities = tracker.find_drifted(namespace=namespace, min_changes=min_changes)

    return {"entities": entities, "count": len(entities)}


@app.get("/api/evolution/viewpoint/{entity_id}")
async def api_viewpoint_evolution(entity_id: str):
    """Get the full viewpoint change history for an entity."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.viewpoint import ViewpointTracker
    tracker = ViewpointTracker(store)
    summary = tracker.get_evolution(entity_id)

    return {
        "entity_id": summary.entity_id,
        "changes": [
            {
                "field": c.field,
                "old_value": c.old_value,
                "new_value": c.new_value,
                "change_type": c.change_type,
                "source": c.source,
                "timestamp": c.timestamp,
            }
            for c in summary.changes
        ],
        "change_count": summary.change_count,
        "latest_change_at": summary.latest_change_at,
    }


# ---------------------------------------------------------------------------
# Phase 4: Synthesis + Obsolescence
# ---------------------------------------------------------------------------


@app.get("/api/synthesis/duplicates")
async def api_synthesis_duplicates(namespace: str = "general", threshold: float = 0.8):
    """Find potential duplicate entities by name similarity."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.synthesis import SynthesisEngine
    engine = SynthesisEngine(store)
    duplicates = engine.find_duplicates(namespace=namespace, similarity_threshold=threshold)

    return {
        "duplicates": [
            {"entity_a": a, "entity_b": b, "similarity": round(s, 3)}
            for a, b, s in duplicates
        ],
        "count": len(duplicates),
    }


@app.post("/api/synthesis/merge")
async def api_synthesis_merge(request: Request):
    """Merge secondary entities into a primary entity."""
    body = await request.json()
    primary_id = body.get("primary_id", "")
    secondary_ids = body.get("secondary_ids", [])

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.synthesis import SynthesisEngine
    engine = SynthesisEngine(store)
    result = engine.merge_entities(primary_id, secondary_ids)

    return {
        "entity_id": result.entity_id,
        "merged_from": result.merged_from,
        "final_confidence": result.final_confidence,
        "property_merges": result.property_merges,
        "conflicts_resolved": result.conflicts_resolved,
    }


@app.post("/api/evolution/obsolescence/mark-outdated")
async def api_obsolescence_mark(request: Request):
    """Mark a specific entity as outdated."""
    body = await request.json()
    entity_id = body.get("entity_id", "")
    reason = body.get("reason", "manual")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.obsolescence import ObsolescenceManager
    mgr = ObsolescenceManager(store)
    return mgr.mark_outdated(entity_id=entity_id, reason=reason)


@app.post("/api/evolution/obsolescence/scan")
async def api_obsolescence_scan(request: Request):
    """Scan for stale entities and mark them outdated."""
    body = await request.json()
    namespace = body.get("namespace", "general")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.obsolescence import ObsolescenceManager
    mgr = ObsolescenceManager(store)
    results = mgr.scan_stale(namespace=namespace)

    return {
        "stale_entities": [
            {
                "entity_id": r.entity_id,
                "action": r.action,
                "reason": r.reason,
                "ttl_days": r.ttl_days,
                "age_days": r.age_days,
            }
            for r in results
        ],
        "count": len(results),
    }


@app.post("/api/evolution/obsolescence/refresh")
async def api_obsolescence_refresh(request: Request):
    """Refresh an entity's timestamp to reset TTL."""
    body = await request.json()
    entity_id = body.get("entity_id", "")

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.evolution.obsolescence import ObsolescenceManager
    mgr = ObsolescenceManager(store)
    return mgr.refresh(entity_id=entity_id)


# ---------------------------------------------------------------------------
# Phase 5: Knowledge Tree + Namespace Configuration
# ---------------------------------------------------------------------------


@app.get("/api/tree")
async def api_tree(namespace: str = "general"):
    """Get the full knowledge tree for a namespace."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.tree import KnowledgeTreeBuilder
    builder = KnowledgeTreeBuilder(store)
    tree = builder.build(namespace=namespace)

    def _serialize_node(node):
        return {
            "entity_id": node.entity_id,
            "name": node.name,
            "type": node.entity_type,
            "importance": node.importance,
            "confidence": node.confidence,
            "is_leaf": node.is_leaf,
            "depth": node.depth,
            "children": [_serialize_node(c) for c in node.children],
        }

    return {"tree": [_serialize_node(n) for n in tree], "namespace": namespace}


@app.get("/api/tree/search")
async def api_tree_search(q: str, namespace: str = "general", limit: int = 20):
    """Search the knowledge tree with breadcrumb context.

    IMPORTANT: Must be registered before /api/tree/{entity_id}/subtree
    to avoid FastAPI route capture.
    """
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.tree import KnowledgeTreeBuilder
    builder = KnowledgeTreeBuilder(store)
    results = builder.search(query=q, namespace=namespace, limit=limit)
    return {"results": results, "count": len(results)}


@app.get("/api/tree/{entity_id}/subtree")
async def api_tree_subtree(entity_id: str, namespace: str = "general"):
    """Get the subtree rooted at a specific entity."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.tree import KnowledgeTreeBuilder
    builder = KnowledgeTreeBuilder(store)
    subtree = builder.get_subtree(entity_id, namespace=namespace)

    if not subtree:
        raise HTTPException(status_code=404, detail="Entity not found")

    def _serialize_node(node):
        return {
            "entity_id": node.entity_id,
            "name": node.name,
            "type": node.entity_type,
            "importance": node.importance,
            "confidence": node.confidence,
            "is_leaf": node.is_leaf,
            "depth": node.depth,
            "children": [_serialize_node(c) for c in node.children],
        }

    return {"subtree": _serialize_node(subtree)}


@app.get("/api/tree/{entity_id}/breadcrumb")
async def api_tree_breadcrumb(entity_id: str, namespace: str = "general"):
    """Get the breadcrumb path from root to a given entity."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.tree import KnowledgeTreeBuilder
    builder = KnowledgeTreeBuilder(store)
    path = builder.get_breadcrumb(entity_id, namespace=namespace)
    return {"path": path}


@app.get("/api/tree/{entity_id}/leaves")
async def api_tree_leaves(entity_id: str, namespace: str = "general"):
    """Get all leaf entities under a given tree node."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.tree import KnowledgeTreeBuilder
    builder = KnowledgeTreeBuilder(store)
    leaves = builder.get_leaves(entity_id, namespace=namespace)
    return {"leaves": leaves}


@app.get("/api/namespaces/config")
async def api_namespace_config_list():
    """List all namespace configurations."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.namespace import NamespaceConfigManager
    mgr = NamespaceConfigManager(store)
    return mgr.list_configs()


@app.get("/api/namespaces/config/{namespace_name}")
async def api_namespace_config_get(namespace_name: str):
    """Get configuration for a specific namespace."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.namespace import NamespaceConfigManager
    mgr = NamespaceConfigManager(store)
    config = mgr.get_config(namespace_name)
    if not config:
        raise HTTPException(status_code=404, detail="Namespace config not found")
    return config


@app.put("/api/namespaces/config/{namespace_name}")
async def api_namespace_config_upsert(namespace_name: str, request: Request):
    """Create or update a namespace configuration."""
    body = await request.json()

    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.namespace import NamespaceConfigManager
    mgr = NamespaceConfigManager(store)
    return mgr.create_or_update(
        namespace_name=namespace_name,
        description=body.get("description", ""),
        retriever_weights=body.get("retriever_weights"),
        entity_type_weights=body.get("entity_type_weights"),
        context_injection=body.get("context_injection"),
    )


@app.delete("/api/namespaces/config/{namespace_name}")
async def api_namespace_config_delete(namespace_name: str):
    """Delete a namespace configuration (does not delete entities)."""
    hm = get_heimdall()
    store = hm.provider.store

    from heimdall.core.namespace import NamespaceConfigManager
    mgr = NamespaceConfigManager(store)
    mgr.delete(namespace_name)
    return {"deleted": True, "namespace": namespace_name}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_entity_type_counts(gs) -> dict:
    counts = {}
    try:
        cur = gs._conn.execute(
            "SELECT entity_type, COUNT(*) FROM heimdall_entities WHERE status='active' GROUP BY entity_type"
        )
        for row in cur:
            counts[row[0]] = row[1]
    except Exception:
        pass
    return counts


def _get_elevator_stats(hm) -> dict:
    try:
        suggestions = hm.provider.get_suggestions()
        return {
            "L1": len(suggestions.get("L1", [])),
            "L2": len(suggestions.get("L2", [])),
            "L3": len(suggestions.get("L3", [])),
        }
    except Exception:
        return {"L1": 0, "L2": 0, "L3": 0}
